"""
Generate a Cold Calling Qualification Questionnaire PPT.
Questions to ask store owners BEFORE visiting their store.
Helps qualify leads and prepare for the pitch.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette ───
BG_DARK       = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD       = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
ACCENT_GOLD   = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_GREEN  = RGBColor(0x00, 0xE6, 0x76)
ACCENT_RED    = RGBColor(0xFF, 0x4D, 0x4D)
ACCENT_BLUE   = RGBColor(0x4D, 0xA8, 0xFF)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xB0, 0xB0, 0xC0)
MEDIUM_GRAY   = RGBColor(0x80, 0x80, 0x99)
SUBTLE_PURPLE = RGBColor(0x6C, 0x5C, 0xE7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=28,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return txBox


def add_rounded_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
    shape.line.width = Pt(1)
    return shape


def add_number_badge(slide, left, top, number, color=ACCENT_ORANGE):
    size = Inches(0.65)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return circle


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)
add_circle(slide, Inches(10.5), Inches(0.5), Inches(2.5), ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(9), Inches(1.2),
             "Store Owner", font_size=52, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.8), Inches(9), Inches(1.2),
             "Qualification Questions", font_size=52, color=ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(1),
             "Ask these questions during your cold call — BEFORE visiting the store.\nThis helps you understand if they are a good fit and prepares you for the meeting.",
             font_size=22, color=LIGHT_GRAY)

add_accent_bar(slide, Inches(1), Inches(6.2), Inches(2), Inches(0.06), ACCENT_BLUE)
add_text_box(slide, Inches(1), Inches(6.4), Inches(6), Inches(0.5),
             "Retail Control Architect™  •  Cold Call Prep Guide",
             font_size=16, color=MEDIUM_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Purpose & How-To-Use
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GOLD)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
             "How to Use This Guide 📋", font_size=40, color=ACCENT_GOLD, bold=True)

instructions = [
    ("📞", "Call the store owner or manager to introduce yourself", ACCENT_BLUE),
    ("🎯", "Use these questions to understand their current situation", ACCENT_ORANGE),
    ("✍️", "Write down their answers — this becomes your prep sheet", ACCENT_GREEN),
    ("🚦", "Green flags = Good fit → Schedule a store visit", ACCENT_GREEN),
    ("🔴", "Red flags = Not ready → Follow up later, don't waste your time", ACCENT_RED),
]

for i, (icon, text, color) in enumerate(instructions):
    y = Inches(1.8) + Inches(i * 1.0)
    add_accent_bar(slide, Inches(1), y, Inches(0.1), Inches(0.75), color)
    add_text_box(slide, Inches(1.5), y + Inches(0.12), Inches(0.5), Inches(0.5),
                 icon, font_size=24)
    add_text_box(slide, Inches(2.2), y + Inches(0.15), Inches(9.5), Inches(0.5),
                 text, font_size=22, color=WHITE)

add_rounded_card(slide, Inches(2.5), Inches(6.5), Inches(7.5), Inches(0.7),
                 RGBColor(0x15, 0x15, 0x2A))
add_text_box(slide, Inches(2.7), Inches(6.55), Inches(7), Inches(0.5),
             "💡 Tip: Keep the conversation friendly. You're helping them, not selling.",
             font_size=16, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Opening Script
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, SUBTLE_PURPLE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
             "Opening Script — What To Say First 🗣️", font_size=38, color=WHITE, bold=True)

# Script card
add_rounded_card(slide, Inches(1), Inches(1.7), Inches(11), Inches(4.5), BG_CARD)

script_lines = [
    ('"Hi [Name], this is Aakash from Retail Control Architect."', WHITE, 22, True),
    ("", WHITE, 10, False),
    ('"I work with clothing & retail store owners to help them', LIGHT_GRAY, 20, False),
    ('stop losing money from stock problems."', LIGHT_GRAY, 20, False),
    ("", WHITE, 10, False),
    ('"I wanted to quickly ask you a few questions to see if', LIGHT_GRAY, 20, False),
    ('I can help your store too. Do you have 3 minutes?"', LIGHT_GRAY, 20, False),
    ("", WHITE, 10, False),
    ("If they say YES → Continue with the questions", ACCENT_GREEN, 18, True),
    ("If they say NO → Ask: 'When would be a better time to call?'", ACCENT_ORANGE, 18, True),
    ("If they say NOT INTERESTED → Thank them and move on", MEDIUM_GRAY, 18, True),
]

y_pos = Inches(1.9)
for text, color, size, bold in script_lines:
    add_text_box(slide, Inches(1.5), y_pos, Inches(10), Inches(0.5),
                 text, font_size=size, color=color, bold=bold)
    y_pos += Inches(0.38)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Basic Store Info Questions
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_number_badge(slide, Inches(1), Inches(0.4), "1", ACCENT_BLUE)
add_text_box(slide, Inches(1.9), Inches(0.4), Inches(8), Inches(0.7),
             "Basic Store Information", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.5),
             "Understanding their business setup", font_size=18, color=MEDIUM_GRAY)

questions = [
    ("Q1", "What type of products do you sell?", "(Clothes, shoes, accessories, etc.)"),
    ("Q2", "How many stores/branches do you have?", "(Single store or multiple locations)"),
    ("Q3", "Roughly how many staff members work in your store?", "(Helps understand team size)"),
    ("Q4", "How long have you been running this business?", "(New vs. experienced owner)"),
]

for i, (q_num, question, note) in enumerate(questions):
    y = Inches(2.0) + Inches(i * 1.25)
    add_rounded_card(slide, Inches(1), y, Inches(10.5), Inches(1.0), BG_CARD)
    add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(1), Inches(0.4),
                 q_num, font_size=18, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.08), Inches(8.5), Inches(0.4),
                 question, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.52), Inches(8.5), Inches(0.4),
                 note, font_size=15, color=MEDIUM_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Stock & Inventory Questions
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)

add_number_badge(slide, Inches(1), Inches(0.4), "2", ACCENT_ORANGE)
add_text_box(slide, Inches(1.9), Inches(0.4), Inches(8), Inches(0.7),
             "Stock & Inventory Situation", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.5),
             "This reveals the REAL pain points", font_size=18, color=MEDIUM_GRAY)

questions = [
    ("Q5", "Do you ever find that your actual stock doesn't match what your records show?",
     "🔥 KEY QUESTION — If YES, they are a perfect fit!"),
    ("Q6", "How do you currently keep track of what comes in and goes out?",
     "(Computer, register, notebook, or just memory?)"),
    ("Q7", "When was the last time you did a full stock count?",
     "(If 'never' or 'long time ago' — very strong lead!)"),
    ("Q8", "Have you ever been surprised by missing items or shortage at month-end?",
     "(This tells you how aware they are of the problem)"),
]

for i, (q_num, question, note) in enumerate(questions):
    y = Inches(2.0) + Inches(i * 1.25)
    add_rounded_card(slide, Inches(1), y, Inches(10.5), Inches(1.0), BG_CARD)
    add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(1), Inches(0.4),
                 q_num, font_size=18, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.08), Inches(8.5), Inches(0.4),
                 question, font_size=19, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.52), Inches(8.5), Inches(0.4),
                 note, font_size=15, color=ACCENT_GOLD)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Staff & Process Questions
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GREEN)

add_number_badge(slide, Inches(1), Inches(0.4), "3", ACCENT_GREEN)
add_text_box(slide, Inches(1.9), Inches(0.4), Inches(8), Inches(0.7),
             "Staff & Daily Process", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.5),
             "Understanding how the store operates day-to-day", font_size=18, color=MEDIUM_GRAY)

questions = [
    ("Q9", "Does each staff member have clear responsibilities — or does everyone do everything?",
     "(Unclear roles = chaos = losses)"),
    ("Q10", "Is there a daily routine or checklist your staff follows for stock?",
     "(If NO — great opportunity to help!)"),
    ("Q11", "Do you feel confident that your staff handles stock carefully?",
     "(Listen for hesitation — that's your signal)"),
    ("Q12", "Have you ever suspected items being taken or mishandled by staff?",
     "(Sensitive question — ask gently)"),
]

for i, (q_num, question, note) in enumerate(questions):
    y = Inches(2.0) + Inches(i * 1.25)
    add_rounded_card(slide, Inches(1), y, Inches(10.5), Inches(1.0), BG_CARD)
    add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(1), Inches(0.4),
                 q_num, font_size=18, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.08), Inches(8.5), Inches(0.4),
                 question, font_size=19, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.52), Inches(8.5), Inches(0.4),
                 note, font_size=15, color=MEDIUM_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Technology & Current Solutions
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, SUBTLE_PURPLE)

add_number_badge(slide, Inches(1), Inches(0.4), "4", SUBTLE_PURPLE)
add_text_box(slide, Inches(1.9), Inches(0.4), Inches(8), Inches(0.7),
             "Current Software & Tools", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.5),
             "Learning what they've already tried", font_size=18, color=MEDIUM_GRAY)

questions = [
    ("Q13", "Do you use any billing software or computer system for your store?",
     "(If YES — which one? If NO — they might be even more open to help)"),
    ("Q14", "Have you ever tried any software specifically for stock management?",
     "(If they tried and failed — they understand the pain deeply)"),
    ("Q15", "How do you check if your store made profit each month?",
     "(Many owners don't know their real profit — big pain point!)"),
]

for i, (q_num, question, note) in enumerate(questions):
    y = Inches(2.0) + Inches(i * 1.5)
    add_rounded_card(slide, Inches(1), y, Inches(10.5), Inches(1.2), BG_CARD)
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(1), Inches(0.4),
                 q_num, font_size=18, color=SUBTLE_PURPLE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.1), Inches(8.5), Inches(0.4),
                 question, font_size=19, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.6), Inches(8.5), Inches(0.4),
                 note, font_size=15, color=MEDIUM_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Money & Willingness Questions
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GOLD)

add_number_badge(slide, Inches(1), Inches(0.4), "5", ACCENT_GOLD)
add_text_box(slide, Inches(1.9), Inches(0.4), Inches(8), Inches(0.7),
             "Willingness & Urgency", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.5),
             "Finding out if they're READY to fix this problem", font_size=18, color=MEDIUM_GRAY)

questions = [
    ("Q16", "If I showed you exactly where you're losing money — would you want to fix it?",
     "🎯 MOST IMPORTANT — If 'YES' → They are QUALIFIED!"),
    ("Q17", "How much do you think stock problems cost your store every month?",
     "(Helps them realize the size of their problem)"),
    ("Q18", "Have you been thinking about improving your stock management recently?",
     "(If YES — they are actively looking for a solution!)"),
    ("Q19", "Would you be open to a free 30-minute store check-up where I visit and diagnose the issues?",
     "🔥 This is your CLOSE — sets up the store visit meeting!"),
]

for i, (q_num, question, note) in enumerate(questions):
    y = Inches(2.0) + Inches(i * 1.25)
    add_rounded_card(slide, Inches(1), y, Inches(10.5), Inches(1.0), BG_CARD)
    add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(1), Inches(0.4),
                 q_num, font_size=18, color=ACCENT_GOLD, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.08), Inches(8.5), Inches(0.4),
                 question, font_size=19, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.52), Inches(8.5), Inches(0.4),
                 note, font_size=15, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Scoring — Green, Yellow, Red Flags
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GREEN)

add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(0.8),
             "How To Score Their Answers 🎯", font_size=38, color=WHITE, bold=True)

# Green flags
add_rounded_card(slide, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.0), RGBColor(0x0A, 0x2A, 0x15))
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(3.3), Inches(0.5),
             "🟢 GREEN — Perfect Fit", font_size=22, color=ACCENT_GREEN, bold=True)
green_flags = [
    "They admit stock doesn't match",
    "They're frustrated with current issues",
    "They've tried other solutions that failed",
    "They said YES to store visit",
    "Multi-store owners",
    "Open to investing in improvement",
]
for i, flag in enumerate(green_flags):
    add_text_box(slide, Inches(1.2), Inches(2.5) + Inches(i * 0.55), Inches(3.1), Inches(0.5),
                 f"✅ {flag}", font_size=15, color=LIGHT_GRAY)

# Yellow flags
add_rounded_card(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.0), RGBColor(0x2A, 0x25, 0x0A))
add_text_box(slide, Inches(5.0), Inches(1.7), Inches(3.3), Inches(0.5),
             "🟡 YELLOW — Follow Up Later", font_size=22, color=ACCENT_GOLD, bold=True)
yellow_flags = [
    "They're 'busy right now'",
    "Not sure if they have a problem",
    "Want to think about it",
    "Asked you to call back later",
    "Interested but hesitant",
    "New store (< 6 months old)",
]
for i, flag in enumerate(yellow_flags):
    add_text_box(slide, Inches(5.2), Inches(2.5) + Inches(i * 0.55), Inches(3.1), Inches(0.5),
                 f"⚠️ {flag}", font_size=15, color=LIGHT_GRAY)

# Red flags
add_rounded_card(slide, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.0), RGBColor(0x2A, 0x15, 0x15))
add_text_box(slide, Inches(9.0), Inches(1.7), Inches(3.3), Inches(0.5),
             "🔴 RED — Not a Fit (Now)", font_size=22, color=ACCENT_RED, bold=True)
red_flags = [
    "Says 'everything is fine'",
    "Doesn't care about stock issues",
    "\"My staff handles everything\"",
    "Refuses to meet or discuss",
    "Only wants free advice",
    "Argumentative / not respectful",
]
for i, flag in enumerate(red_flags):
    add_text_box(slide, Inches(9.2), Inches(2.5) + Inches(i * 0.55), Inches(3.1), Inches(0.5),
                 f"❌ {flag}", font_size=15, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Closing Script — Booking the Visit
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
             "Closing Script — Book the Store Visit 📅", font_size=38, color=WHITE, bold=True)

# Closing script card
add_rounded_card(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.0), BG_CARD)

close_lines = [
    ("Based on what you have told me, I think I can really help your store.", WHITE, 22, True),
    ("", WHITE, 8, False),
    ("Here is what I would like to do -- completely FREE, no obligation:", ACCENT_GOLD, 20, True),
    ("", WHITE, 8, False),
    ("  I will visit your store for 30 minutes", LIGHT_GRAY, 20, False),
    ("  Walk through your stock process with you", LIGHT_GRAY, 20, False),
    ("  Show you exactly where money might be leaking", LIGHT_GRAY, 20, False),
    ("  Give you a clear picture of what can be improved", LIGHT_GRAY, 20, False),
    ("", WHITE, 8, False),
    ("Would Tuesday or Thursday work better for you?", ACCENT_GREEN, 22, True),
    ("(Always give 2 options -- makes it easier for them to choose)", MEDIUM_GRAY, 15, False),
    ("", WHITE, 8, False),
    ("If they hesitate:", ACCENT_ORANGE, 18, True),
    ("No worries! There is absolutely no cost and no pressure. I just want to show you", LIGHT_GRAY, 17, False),
    ("what I found in other stores like yours -- it could save you lakhs.", LIGHT_GRAY, 17, False),
]

y_pos = Inches(1.85)
for text, color, size, bold in close_lines:
    add_text_box(slide, Inches(1.5), y_pos, Inches(10), Inches(0.5),
                 text, font_size=size, color=color, bold=bold)
    y_pos += Inches(0.35)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: Notes template — What to write down
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(0.8),
             "Notes Template — Fill This For Each Call 📝", font_size=38, color=WHITE, bold=True)

# Template card
fields = [
    ("Store Name:", "___________________________________"),
    ("Owner Name:", "___________________________________"),
    ("Phone Number:", "___________________________________"),
    ("Store Type:", "Clothing / Footwear / Accessories / Other: _______"),
    ("Number of Stores:", "___________________________________"),
    ("Staff Count:", "___________________________________"),
    ("Current Software:", "None / Billing Only / Stock Mgmt: ____________"),
    ("Main Problem:", "___________________________________"),
    ("Interested in Visit?", "YES  /  MAYBE  /  NO"),
    ("Score:", "🟢 GREEN  /  🟡 YELLOW  /  🔴 RED"),
    ("Visit Date/Time:", "___________________________________"),
    ("Extra Notes:", "___________________________________"),
]

for i, (label, value) in enumerate(fields):
    y = Inches(1.5) + Inches(i * 0.48)
    add_text_box(slide, Inches(1.2), y, Inches(3), Inches(0.4),
                 label, font_size=16, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(4.2), y, Inches(8), Inches(0.4),
                 value, font_size=16, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: Quick Do's & Don'ts
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GOLD)

add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(0.8),
             "Cold Call Do's & Don'ts 📌", font_size=38, color=ACCENT_GOLD, bold=True)

# DO's
add_rounded_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2), RGBColor(0x0A, 0x2A, 0x15))
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.5),
             "✅ DO's", font_size=28, color=ACCENT_GREEN, bold=True)

dos = [
    "Be friendly and warm — smile while talking",
    "Listen MORE than you talk (80/20 rule)",
    "Ask open-ended questions",
    "Take notes during the call",
    "Offer genuine help, not a sales pitch",
    "Respect their time — keep it under 5 min",
    "Confirm the meeting before hanging up",
    "Follow up if they asked you to call back",
]

for i, text in enumerate(dos):
    add_text_box(slide, Inches(1.2), Inches(2.5) + Inches(i * 0.5), Inches(4.8), Inches(0.5),
                 f"✅  {text}", font_size=15, color=LIGHT_GRAY)

# DON'Ts
add_rounded_card(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.2), RGBColor(0x2A, 0x15, 0x15))
add_text_box(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.5),
             "❌ DON'Ts", font_size=28, color=ACCENT_RED, bold=True)

donts = [
    "Don't start with pricing or cost",
    "Don't criticize their current way of working",
    "Don't use complicated words or tech jargon",
    "Don't argue if they say no",
    "Don't call during busy hours (11am-1pm)",
    "Don't promise unrealistic results",
    "Don't sound like a telemarketer",
    "Don't rush — let them think and respond",
]

for i, text in enumerate(donts):
    add_text_box(slide, Inches(7.2), Inches(2.5) + Inches(i * 0.5), Inches(4.8), Inches(0.5),
                 f"❌  {text}", font_size=15, color=LIGHT_GRAY)


# ─── Save ───
output_path = "/Users/aakash/Desktop/Week_3/Cold_Call_Qualification_Questions.pptx"
prs.save(output_path)
print(f"✅ Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
