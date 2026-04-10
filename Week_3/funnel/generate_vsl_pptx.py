#!/usr/bin/env python3
"""
Generate Retail Control Architect™ — Complete VSL PowerPoint Presentation
70+ slides with animations, transitions, and professional design.
Ready to record as a Video Sales Letter.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import os

# ── Colour Palette ──────────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x0A, 0x0E, 0x1A)
DEEP_BLUE   = RGBColor(0x0F, 0x1B, 0x2D)
MID_BLUE    = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_GOLD = RGBColor(0xD4, 0xA0, 0x1E)
ACCENT_TEAL = RGBColor(0x17, 0xA2, 0xB8)
WARM_RED    = RGBColor(0xE7, 0x4C, 0x3C)
SUCCESS_GRN = RGBColor(0x2E, 0xCC, 0x71)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
SOFT_WHITE  = RGBColor(0xF0, 0xF0, 0xF0)
ORANGE      = RGBColor(0xF3, 0x9C, 0x12)
PURPLE      = RGBColor(0x9B, 0x59, 0xB6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ── Animation Helpers ───────────────────────────────────────────────

def add_fade_in(slide, shape, delay_ms=0, duration_ms=800):
    """Add a fade-in entrance animation to a shape."""
    timing = slide.element.find(qn('p:timing'))
    if timing is None:
        timing = etree.SubElement(slide.element, qn('p:timing'))
    
    tn_lst = timing.find(qn('p:tnLst'))
    if tn_lst is None:
        tn_lst = etree.SubElement(timing, qn('p:tnLst'))
    
    par = tn_lst.find(qn('p:par'))
    if par is None:
        par = etree.SubElement(tn_lst, qn('p:par'))
        c_tn_root = etree.SubElement(par, qn('p:cTn'), attrib={
            'id': '1', 'dur': 'indefinite', 'restart': 'never', 'nodeType': 'tmRoot'
        })
        child_tn_lst = etree.SubElement(c_tn_root, qn('p:childTnLst'))
    else:
        c_tn_root = par.find(qn('p:cTn'))
        child_tn_lst = c_tn_root.find(qn('p:childTnLst'))
    
    # Get next id
    all_ids = [int(e.get('id', '0')) for e in slide.element.iter() if e.get('id')]
    next_id = max(all_ids) + 1 if all_ids else 2
    
    seq = child_tn_lst.find(qn('p:seq'))
    if seq is None:
        seq_par = etree.SubElement(child_tn_lst, qn('p:seq'), attrib={
            'concurrent': '1', 'nextAc': 'seek'
        })
        seq_ctn = etree.SubElement(seq_par, qn('p:cTn'), attrib={
            'id': str(next_id), 'dur': 'indefinite', 'nodeType': 'mainSeq'
        })
        next_id += 1
        seq_child = etree.SubElement(seq_ctn, qn('p:childTnLst'))
    else:
        seq_ctn = seq.find(qn('p:cTn'))
        seq_child = seq_ctn.find(qn('p:childTnLst'))
        if seq_child is None:
            seq_child = etree.SubElement(seq_ctn, qn('p:childTnLst'))
    
    # Create animation
    anim_par = etree.SubElement(seq_child, qn('p:par'))
    anim_ctn = etree.SubElement(anim_par, qn('p:cTn'), attrib={
        'id': str(next_id), 'fill': 'hold'
    })
    next_id += 1
    
    stCondLst = etree.SubElement(anim_ctn, qn('p:stCondLst'))
    cond = etree.SubElement(stCondLst, qn('p:cond'), attrib={'delay': str(delay_ms)})
    
    anim_child = etree.SubElement(anim_ctn, qn('p:childTnLst'))
    
    inner_par = etree.SubElement(anim_child, qn('p:par'))
    inner_ctn = etree.SubElement(inner_par, qn('p:cTn'), attrib={
        'id': str(next_id), 'presetID': '10', 'presetClass': 'entr',
        'presetSubtype': '0', 'fill': 'hold', 'nodeType': 'withEffect'
    })
    next_id += 1
    
    inner_stCond = etree.SubElement(inner_ctn, qn('p:stCondLst'))
    etree.SubElement(inner_stCond, qn('p:cond'), attrib={'delay': '0'})
    
    inner_child = etree.SubElement(inner_ctn, qn('p:childTnLst'))
    
    # Alpha animation (fade)
    anim_effect = etree.SubElement(inner_child, qn('p:animEffect'), attrib={
        'transition': 'in', 'filter': 'fade'
    })
    effect_ctn = etree.SubElement(anim_effect, qn('p:cBhvr'))
    effect_inner_ctn = etree.SubElement(effect_ctn, qn('p:cTn'), attrib={
        'id': str(next_id), 'dur': str(duration_ms)
    })
    next_id += 1
    tgt_el = etree.SubElement(effect_ctn, qn('p:tgtEl'))
    sp_tgt = etree.SubElement(tgt_el, qn('p:spTgt'), attrib={
        'spid': str(shape.shape_id)
    })


def add_slide_transition(slide, trans_type='fade', duration_ms=700, advance_ms=None):
    """Add a slide transition."""
    transition = etree.SubElement(slide.element, qn('p:transition'), attrib={
        'spd': 'med'
    })
    if advance_ms:
        transition.set('advTm', str(advance_ms))
    
    if trans_type == 'fade':
        etree.SubElement(transition, qn('p:fade'))
    elif trans_type == 'push':
        etree.SubElement(transition, qn('p:push'))
    elif trans_type == 'wipe':
        etree.SubElement(transition, qn('p:wipe'))
    elif trans_type == 'cover':
        etree.SubElement(transition, qn('p:cover'))
    elif trans_type == 'split':
        etree.SubElement(transition, qn('p:split'))


# ── Slide Builders ──────────────────────────────────────────────────

def set_bg(slide, r, g, b):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_shape_with_gradient(slide, left, top, width, height, color1, color2):
    """Add a rectangle with gradient fill."""
    shape = slide.shapes.add_shape(
        1, left, top, width, height  # MSO_SHAPE.RECTANGLE = 1
    )
    shape.fill.gradient()
    shape.fill.gradient_stops[0].color.rgb = color1
    shape.fill.gradient_stops[0].position = 0.0
    shape.fill.gradient_stops[1].color.rgb = color2
    shape.fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=32,
                 color=WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 font_name='Calibri'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, 
                      default_size=28, default_color=WHITE, default_bold=False,
                      alignment=PP_ALIGN.CENTER):
    """Add text box with multiple formatted lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, color, bold = line_data, default_size, default_color, default_bold
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            color = line_data[2] if len(line_data) > 2 else default_color
            bold = line_data[3] if len(line_data) > 3 else default_bold
        
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Calibri'
        p.alignment = alignment
        p.space_after = Pt(8)
    
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_GOLD):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(1, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def make_dark_slide(text, subtitle=None, accent_color=ACCENT_GOLD,
                    font_size=40, trans='fade', sub_size=24):
    """Standard dark slide with centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(slide, 0x0A, 0x0E, 0x1A)
    
    # Accent line top
    add_accent_line(slide, Inches(4), Inches(0.3), Inches(5.333), accent_color)
    
    # Main text
    main = add_text_box(
        slide, Inches(1.5), Inches(2.0), Inches(10.333), Inches(3.0),
        text, font_size=font_size, color=WHITE, bold=True
    )
    
    if subtitle:
        sub = add_text_box(
            slide, Inches(2.0), Inches(5.2), Inches(9.333), Inches(1.5),
            subtitle, font_size=sub_size, color=LIGHT_GRAY, bold=False
        )
        add_fade_in(slide, sub, delay_ms=600)
    
    # Accent line bottom
    add_accent_line(slide, Inches(4), Inches(7.0), Inches(5.333), accent_color)
    
    add_fade_in(slide, main, delay_ms=0)
    add_slide_transition(slide, trans)
    return slide


def make_impact_slide(text, highlight_color=ACCENT_GOLD, font_size=48, trans='fade'):
    """High-impact single-statement slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0x0A, 0x0E, 0x1A)
    
    main = add_text_box(
        slide, Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0),
        text, font_size=font_size, color=highlight_color, bold=True
    )
    
    add_fade_in(slide, main, delay_ms=0, duration_ms=1000)
    add_slide_transition(slide, trans)
    return slide


def make_proof_slide(headline, bullets, trans='fade'):
    """Slide with headline + bullet points for proof/results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0x0A, 0x0E, 0x1A)
    
    # Header bar
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MID_BLUE
    bar.line.fill.background()
    
    add_text_box(
        slide, Inches(0.5), Inches(0.15), Inches(12.333), Inches(0.9),
        headline, font_size=32, color=ACCENT_GOLD, bold=True
    )
    
    # Bullets
    lines = [(b, 24, WHITE, False) for b in bullets]
    box = add_multiline_box(
        slide, Inches(1.5), Inches(1.8), Inches(10.333), Inches(5.0),
        lines, alignment=PP_ALIGN.LEFT
    )
    
    add_fade_in(slide, box, delay_ms=400)
    add_slide_transition(slide, trans)
    return slide


def make_phase_title(phase_num, title, emoji="", color=ACCENT_GOLD):
    """Phase title/section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 0x0A, 0x0E, 0x1A)
    
    # Top accent
    bar = slide.shapes.add_shape(1, 0, Inches(2.5), SLIDE_W, Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MID_BLUE
    bar.line.fill.background()
    
    phase_label = add_text_box(
        slide, Inches(1), Inches(2.6), Inches(11.333), Inches(0.8),
        f"PHASE {phase_num}", font_size=20, color=color, bold=True
    )
    
    main = add_text_box(
        slide, Inches(1), Inches(3.2), Inches(11.333), Inches(1.5),
        f"{emoji}  {title}" if emoji else title,
        font_size=44, color=WHITE, bold=True
    )
    
    add_accent_line(slide, Inches(5), Inches(4.8), Inches(3.333), color)
    
    add_fade_in(slide, phase_label, delay_ms=0)
    add_fade_in(slide, main, delay_ms=300, duration_ms=1000)
    add_slide_transition(slide, 'fade')
    return slide


# ════════════════════════════════════════════════════════════════════
#                     BUILD THE PRESENTATION
# ════════════════════════════════════════════════════════════════════

# ── SLIDE 0: TITLE ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)

add_accent_line(slide, Inches(3.5), Inches(1.5), Inches(6.333))
t1 = add_text_box(slide, Inches(1), Inches(1.8), Inches(11.333), Inches(1.2),
    "RETAIL CONTROL ARCHITECT™", font_size=42, color=ACCENT_GOLD)
t2 = add_text_box(slide, Inches(1), Inches(3.0), Inches(11.333), Inches(1.0),
    "Video Sales Letter", font_size=28, color=WHITE)
t3 = add_text_box(slide, Inches(2), Inches(4.2), Inches(9.333), Inches(1.2),
    "Stop the Silent Bleed in Your Retail Store", font_size=22, color=LIGHT_GRAY, bold=False)
add_accent_line(slide, Inches(3.5), Inches(5.8), Inches(6.333))
add_fade_in(slide, t1, delay_ms=0, duration_ms=1200)
add_fade_in(slide, t2, delay_ms=600)
add_fade_in(slide, t3, delay_ms=1200)
add_slide_transition(slide, 'fade')


# ════════════════════════════════════════════════════════════════════
# PHASE 1: PATTERN INTERRUPT (Slides 1-3)
# ════════════════════════════════════════════════════════════════════
make_phase_title(1, "PATTERN INTERRUPT", "⚡")

# Slide 1
make_impact_slide(
    "There is a thief inside your store\nright now.",
    highlight_color=WARM_RED, font_size=52
)

# Slide 2
make_dark_slide(
    "But here's the thing…\n\n"
    "You hired him.\nYou pay his salary.\n"
    "And he steals from you\nevery single day.",
    font_size=36
)

# Slide 3
make_dark_slide(
    "In the next few minutes, I'm going to show you\n"
    "how this 'invisible thief' is silently draining\n"
    "₹2 to 5 Lakhs from your retail store every year…",
    subtitle="And exactly how to catch him — in 30 days flat.",
    accent_color=WARM_RED, font_size=34
)


# ════════════════════════════════════════════════════════════════════
# PHASE 2: BIG PROMISE + INTRO (Slides 4-6)
# ════════════════════════════════════════════════════════════════════
make_phase_title(2, "THE BIG PROMISE", "🎯")

# Slide 4 — Big Promise
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)
box = add_multiline_box(
    slide, Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5),
    [
        ("How to stop losing ₹40,000+ every single month", 36, ACCENT_GOLD, True),
        ("from your clothing store…", 36, ACCENT_GOLD, True),
        ("", 16, WHITE, False),
        ("Without buying expensive software", 26, WHITE, False),
        ("Without firing your staff", 26, WHITE, False),
        ("Without spending hours reconciling stock every night", 26, WHITE, False),
        ("", 16, WHITE, False),
        ("Even if you've tried billing software before", 22, LIGHT_GRAY, False),
        ("and it didn't work.", 22, LIGHT_GRAY, False),
    ]
)
add_fade_in(slide, box, delay_ms=0, duration_ms=1000)
add_slide_transition(slide, 'fade')

# Slide 5
make_dark_slide("Hey, I'm Aakash Savant.", font_size=44)

# Slide 6
make_impact_slide("And let me ask you this…",
    highlight_color=WHITE, font_size=44)


# ════════════════════════════════════════════════════════════════════
# PHASE 3: YES LADDER (Slides 7-13)
# ════════════════════════════════════════════════════════════════════
make_phase_title(3, "THE YES LADDER", "😤")

make_dark_slide(
    "Are you sick and tired of your stock count\n"
    "NEVER matching what your\nbilling software says?",
    accent_color=WARM_RED, font_size=36
)

make_dark_slide(
    "Do you wish you could stop spending\n"
    "2 hours every night trying to find\n"
    "where items went missing?",
    font_size=36
)

make_dark_slide(
    "Are you looking for a way to finally\n"
    "run your store without depending on\n"
    "one 'trusted' staff member who may\n"
    "or may not be stealing from you?",
    font_size=34
)

make_dark_slide(
    "Or do you simply want to go home on time,\n"
    "knowing that every single item in your store\n"
    "is exactly where it should be?",
    font_size=34
)

make_impact_slide(
    "If this sounds like you…\nthen you're in luck.",
    highlight_color=SUCCESS_GRN, font_size=44
)

make_dark_slide(
    "Because in the next few minutes,\n"
    "I'm going to reveal the single most effective\n"
    "stock control breakthrough\n"
    "in the Indian retail industry.",
    font_size=32
)

make_impact_slide("Sound good?\n\nExcellent!",
    highlight_color=ACCENT_GOLD, font_size=48)


# ════════════════════════════════════════════════════════════════════
# PHASE 4: EXCLUSIVITY + PROOF (Slides 14-20)
# ════════════════════════════════════════════════════════════════════
make_phase_title(4, "EXCLUSIVITY + PROOF BLAST", "🏆")

make_dark_slide("Now before I reveal how this works…", font_size=38)

make_impact_slide(
    "This is NOT for everybody.",
    highlight_color=WARM_RED, font_size=52
)

make_dark_slide(
    "I only take on 3 new stores per month.\n\n"
    "For reasons you'll understand\n"
    "in just a moment…",
    subtitle="I have to be extremely selective about who I work with.",
    font_size=34
)

make_impact_slide(
    "But first, let me SHOW you\nwhat I'm talking about…",
    highlight_color=ACCENT_GOLD, font_size=44
)

# Proof slide 1
make_proof_slide("📸  REAL RESULTS FROM REAL STORES", [
    "🏪  Store in Pune — Mismatch dropped from 17% to under 2% in 45 days",
    "",
    "💰  ₹3.2 Lakhs recovered in hidden stock losses",
    "",
    "📊  847 mismatched items → 12 items. Stock accuracy: 99.4%",
    "",
    '💬  Owner: "I sleep better at night knowing my stock is right."',
])

# Proof slide 2
make_proof_slide("📸  MORE RESULTS — DIFFERENT STORES, SAME SYSTEM", [
    "🏬  Dead stock worth ₹4.2 Lakhs identified and liquidated in 3 weeks",
    "",
    "⏱️  Nightly reconciliation: 2 hours → 10 minutes",
    "",
    "🏪  One owner opened his second branch within 6 months",
    "     — his operations were that clean",
    "",
    '💬  "The system runs even when I\'m not in the store."',
])


# ════════════════════════════════════════════════════════════════════
# PHASE 5: CREDIBILITY + ORIGIN STORY (Slides 21-30)
# ════════════════════════════════════════════════════════════════════
make_phase_title(5, "THE ORIGIN STORY", "👤")

make_dark_slide(
    "Now before I tell you exactly\nhow we achieve these results…",
    subtitle="Let me briefly explain who I am and\nwhy you should listen to me at all.",
    font_size=36
)

make_dark_slide(
    "My name is Aakash Savant\n\n"
    "Founder of Retail Control Architect™",
    subtitle="A specialized stock control system designed exclusively\n"
    "for owner-operated clothing stores in India.",
    font_size=36, sub_size=22
)

make_dark_slide(
    "We help retail store owners like you\n"
    "eliminate stock mismatch,\n"
    "recover hidden losses,\n"
    "and build a store that runs on a SYSTEM.",
    font_size=34
)

make_impact_slide("Here's what happened…",
    highlight_color=ACCENT_GOLD, font_size=48)

# The story
make_dark_slide(
    "I walked into a store and asked\none simple question:\n\n"
    "\"How many units of your\nbest-selling shirt do you have?\"",
    accent_color=ACCENT_TEAL, font_size=34
)

# The reveal
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)
box = add_multiline_box(
    slide, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.5),
    [
        ("The screen said:", 28, LIGHT_GRAY, False),
        ("34 units", 56, SUCCESS_GRN, True),
        ("", 16, WHITE, False),
        ("We went to the rack. Counted.", 28, LIGHT_GRAY, False),
        ("21 units", 56, WARM_RED, True),
        ("", 16, WHITE, False),
        ("13 shirts — gone.", 36, WHITE, True),
        ("No record. No explanation. Just… gone.", 24, LIGHT_GRAY, False),
    ]
)
add_fade_in(slide, box, delay_ms=0, duration_ms=1000)
add_slide_transition(slide, 'fade')

make_dark_slide(
    "And that was just ONE item\nout of 2,000+.\n\n"
    "The store was losing\n₹4 Lakhs a year.",
    accent_color=WARM_RED, font_size=36
)

# The big insight
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)
box = add_multiline_box(
    slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(5.0),
    [
        ("This is NOT a software problem.", 36, WARM_RED, True),
        ("❌", 28, WARM_RED, False),
        ("", 12, WHITE, False),
        ("It's NOT a staff problem.", 36, WARM_RED, True),
        ("❌", 28, WARM_RED, False),
        ("", 12, WHITE, False),
        ("It is a SYSTEM problem.", 42, SUCCESS_GRN, True),
        ("✅", 32, SUCCESS_GRN, False),
    ]
)
add_fade_in(slide, box, delay_ms=0, duration_ms=1200)
add_slide_transition(slide, 'fade')

make_dark_slide(
    "So I built one.\n\n"
    "A simple 5-step system that any store owner\n"
    "can implement in 30 days.",
    subtitle="Studied the best inventory systems worldwide.\nDistilled it into one framework.",
    font_size=34
)


# ════════════════════════════════════════════════════════════════════
# PHASE 6: CLOSE THE LOOP + SYSTEM REVEAL (Slides 31-38)
# ════════════════════════════════════════════════════════════════════
make_phase_title(6, "THE SYSTEM REVEALED", "🔓")

make_impact_slide(
    "Remember that\ninvisible thief?",
    highlight_color=WARM_RED, font_size=52
)

# Loop close
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)
box = add_multiline_box(
    slide, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.5),
    [
        ("The thief is not a person.", 32, WHITE, False),
        ("It's not your staff member.", 32, WHITE, False),
        ("It's not your supplier.", 32, WHITE, False),
        ("", 16, WHITE, False),
        ("The invisible thief is your", 32, LIGHT_GRAY, False),
        ("BROKEN SYSTEM", 52, WARM_RED, True),
        ("", 12, WHITE, False),
        ("And you've been paying for it", 26, LIGHT_GRAY, False),
        ("every single day.", 26, LIGHT_GRAY, False),
    ]
)
add_fade_in(slide, box, delay_ms=0, duration_ms=1200)
add_slide_transition(slide, 'fade')

make_impact_slide(
    "This is NOT like anything\nyou've heard before.",
    highlight_color=ACCENT_GOLD, font_size=44
)

make_dark_slide(
    "Pay very close attention —\n"
    "this could be the breakthrough\nyour store has been waiting for.",
    font_size=36
)

# System reveal
make_proof_slide("⚙️  THE 5-STEP RETAIL CONTROL ARCHITECT™ SYSTEM", [
    "Step 1 ▸ FIND THE LEAKS — Full forensic audit of real vs system stock",
    "",
    "Step 2 ▸ ORGANIZE EVERY ITEM — Unique tagging for every SKU",
    "",
    "Step 3 ▸ STAFF ACCOUNTABILITY — Daily checklists per staff member",
    "",
    "Step 4 ▸ WEEKLY TRACKING — One-page report, 10 minutes",
    "",
    "Step 5 ▸ LOCK THE PROCESS — Written rules that run without you",
])

make_dark_slide(
    "This is what we do\nday-in, day-out\nat Retail Control Architect™.",
    font_size=36
)

make_dark_slide(
    "And we've done it for store owners\njust like you…",
    font_size=38
)


# ════════════════════════════════════════════════════════════════════
# PHASE 7: MORE PROOF (Slides 39-44)
# ════════════════════════════════════════════════════════════════════
make_phase_title(7, "UNDENIABLE PROOF", "📸")

make_proof_slide("💬  WHAT STORE OWNERS ARE SAYING", [
    '🏪 Rajesh, Pune:',
    '   "My stock matches exactly now. I sleep better at night."',
    '',
    '🏬 Meera, Bengaluru:',
    '   "We found ₹1.8 Lakhs in dead stock we didn\'t know existed."',
    '',
    '🏪 Vikram, Mumbai:',
    '   "My staff actually follows a process now. It\'s like a different store."',
])

make_dark_slide(
    "We find and eliminate hidden stock losses\n"
    "for retail store owners every single month.\n\n"
    "Using a system that is simple,\nrepeatable, and guaranteed to work.",
    font_size=32
)

make_dark_slide(
    "We do all the detective work for you —\n"
    "the forensic audit, the SKU cleanup, the staff training —\n\n"
    "so you don't have to become a tech expert\nor an operations genius.",
    font_size=30
)

# Dream picture
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)
box = add_multiline_box(
    slide, Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5),
    [
        ("Imagine…", 40, ACCENT_GOLD, True),
        ("", 12, WHITE, False),
        ("Walking on the shop floor, knowing that every item", 28, WHITE, False),
        ("on every rack is exactly where your system says it is.", 28, WHITE, False),
        ("", 16, WHITE, False),
        ("Leaving the store on time, with a clean daily report,", 28, WHITE, False),
        ("and zero surprises waiting for you tomorrow.", 28, WHITE, False),
        ("", 16, WHITE, False),
        ("That's what we deliver. Guaranteed.", 32, SUCCESS_GRN, True),
    ]
)
add_fade_in(slide, box, delay_ms=0, duration_ms=1000)
add_slide_transition(slide, 'fade')


# ════════════════════════════════════════════════════════════════════
# PHASE 8: OBJECTION HANDLING (Slides 45-52)
# ════════════════════════════════════════════════════════════════════
make_phase_title(8, "OBJECTION HANDLING", "🛡️")

make_dark_slide(
    "And by the way…\nyou might be thinking…",
    font_size=40
)

# Objection 1
make_impact_slide(
    '"But Aakash, I already use\nbilling software.\nDoesn\'t that handle inventory?"',
    highlight_color=LIGHT_GRAY, font_size=36
)

make_dark_slide(
    "I don't blame you.\nI used to think the same way.\n\n"
    "Until I discovered that billing software\nonly RECORDS data.\n\n"
    "It does NOT create the daily habits\nthat PREVENT stock from going missing.",
    font_size=30
)

# Objection 2
make_impact_slide(
    '"My staff won\'t follow\nany new system."',
    highlight_color=LIGHT_GRAY, font_size=40
)

make_dark_slide(
    "That's exactly why this system has\n"
    "built-in daily checks that\nphysically CANNOT be skipped.\n\n"
    "We don't hope for compliance.\n"
    "We engineer it.",
    font_size=32, accent_color=SUCCESS_GRN
)

# Objection 3
make_impact_slide(
    '"My store is too small for this."',
    highlight_color=LIGHT_GRAY, font_size=42
)

make_dark_slide(
    "Smaller stores lose a HIGHER percentage\n"
    "of revenue to stock mismatch.\n\n"
    "₹2,000/day × 365 = ₹7.3 Lakhs/year.\n\n"
    "This system was BUILT\nfor stores like yours.",
    font_size=32, accent_color=WARM_RED
)

make_impact_slide(
    "Are you beginning to see\nhow easy it can be?",
    highlight_color=SUCCESS_GRN, font_size=42
)


# ════════════════════════════════════════════════════════════════════
# PHASE 9: OFFER + GUARANTEE (Slides 53-60)
# ════════════════════════════════════════════════════════════════════
make_phase_title(9, "THE OFFER + GUARANTEE", "💎")

make_proof_slide("✨  THE RETAIL CONTROL ARCHITECT™ SYSTEM IS…", [
    "🆕  NEW — No one else installs a stock PROCESS in your store",
    "",
    "✅  SIMPLE — Takes 10 minutes a day to maintain",
    "",
    "🤝  CONVENIENT — We do the heavy lifting for you",
    "",
    "📊  PREDICTABLE — Measurable mismatch reduction in 30 days",
    "",
    "🏆  And best of all…",
])

make_impact_slide(
    "My team and I would love to\nfind YOUR biggest stock leak\n\n— FOR FREE.",
    highlight_color=SUCCESS_GRN, font_size=44
)

make_impact_slide(
    "And I'm so confident we can\nthat I'll even guarantee it.",
    highlight_color=ACCENT_GOLD, font_size=42
)

# Guarantee slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)

# Green border box
bar = slide.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
bar.fill.solid()
bar.fill.fore_color.rgb = RGBColor(0x0D, 0x2B, 0x1A)
bar.line.color.rgb = SUCCESS_GRN
bar.line.width = Pt(3)

box = add_multiline_box(
    slide, Inches(2.0), Inches(1.8), Inches(9.333), Inches(4.0),
    [
        ("🔒  THE STOCK CERTAINTY GUARANTEE™", 32, SUCCESS_GRN, True),
        ("", 16, WHITE, False),
        ("If after 30 days, your stock mismatch", 28, WHITE, False),
        ("does not reduce measurably,", 28, WHITE, False),
        ("", 10, WHITE, False),
        ("I will keep working with you", 32, ACCENT_GOLD, True),
        ("for FREE until it does.", 32, ACCENT_GOLD, True),
        ("", 16, WHITE, False),
        ("We don't want your money", 24, LIGHT_GRAY, False),
        ("if your stock doesn't match.", 24, LIGHT_GRAY, False),
    ]
)
add_fade_in(slide, box, delay_ms=0, duration_ms=1200)
add_slide_transition(slide, 'fade')

make_impact_slide(
    "Want a completely risk-free way\nto find out how much\nyour store is losing?",
    highlight_color=WHITE, font_size=40
)

# Free offer
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)
box = add_multiline_box(
    slide, Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5),
    [
        ("For a limited time, claim your", 28, WHITE, False),
        ("", 8, WHITE, False),
        ("FREE 30-Minute", 48, SUCCESS_GRN, True),
        ("Store Check-Up", 48, SUCCESS_GRN, True),
        ("", 16, WHITE, False),
        ("We'll find your #1 leak source together", 26, LIGHT_GRAY, False),
        ("and tell you exactly how much it costs you", 26, LIGHT_GRAY, False),
        ("every single month.", 26, LIGHT_GRAY, False),
    ]
)
add_fade_in(slide, box, delay_ms=0, duration_ms=1000)
add_slide_transition(slide, 'fade')

make_proof_slide("📋  ON THIS FREE CALL, WE'LL COVER:", [
    "✅  How you can reduce stock mismatch to under 2% in 30 days",
    "",
    "✅  The #1 thing you should NEVER do when fixing inventory",
    "",
    "✅  3 common mistakes clothing store owners make",
    "     that silently drain their profits",
    "",
    "🚫  No cost. No obligation. No pressure.",
])


# ════════════════════════════════════════════════════════════════════
# PHASE 10: CTA + URGENCY CLOSE (Slides 61-70)
# ════════════════════════════════════════════════════════════════════
make_phase_title(10, "THE CLOSE", "🔥")

make_impact_slide("But a word of warning…",
    highlight_color=WARM_RED, font_size=48)

make_dark_slide(
    "This is NOT some quick-fix magic solution.\n\n"
    "This is for legitimate, hard-working\nretail store owners who are ready to\n"
    "invest in fixing their operations.",
    font_size=30
)

make_dark_slide(
    "You must be willing to commit to the process.\n\n"
    "I have a strict 'Single System Rule' —\n"
    "we eliminate all parallel diaries\nand manual shortcuts.",
    font_size=30
)

make_impact_slide(
    "If you can agree to that…\n\n"
    "this is truly the opportunity\nof a lifetime.",
    highlight_color=ACCENT_GOLD, font_size=40
)

make_dark_slide(
    "This call is NOT a sales call.\n\n"
    "It's a free diagnostic session\n"
    "to find your biggest leak\nand see if we're a good fit.",
    font_size=32
)

make_dark_slide(
    "But this video is being seen by\nhundreds of store owners.\n\n"
    "I take on only 3 new stores per month.\n\n"
    "First come, first served.",
    font_size=32, accent_color=WARM_RED
)

# CTA slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)

# CTA Button
btn = slide.shapes.add_shape(
    5, Inches(3.0), Inches(4.5), Inches(7.333), Inches(1.2)  # Rounded rect
)
btn.fill.solid()
btn.fill.fore_color.rgb = SUCCESS_GRN
btn.line.fill.background()

add_text_box(slide, Inches(3.0), Inches(4.6), Inches(7.333), Inches(1.0),
    "👉  BOOK MY FREE STORE CHECK-UP  →", font_size=28, color=WHITE)

box = add_multiline_box(
    slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(2.5),
    [
        ("Don't wait.", 40, WHITE, True),
        ("Click the button below NOW.", 36, ACCENT_GOLD, True),
        ("Book the next available slot.", 28, LIGHT_GRAY, False),
    ]
)
add_fade_in(slide, btn, delay_ms=0, duration_ms=800)
add_fade_in(slide, box, delay_ms=400)
add_slide_transition(slide, 'fade')

# Two options
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11.333), Inches(0.8),
    "You've got two options here…", font_size=32, color=ACCENT_GOLD)

# Option 1 box
o1 = slide.shapes.add_shape(1, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
o1.fill.solid()
o1.fill.fore_color.rgb = RGBColor(0x2D, 0x0A, 0x0A)
o1.line.color.rgb = WARM_RED
o1.line.width = Pt(2)

add_multiline_box(
    slide, Inches(1.2), Inches(2.0), Inches(4.8), Inches(4.0),
    [
        ("OPTION 1", 24, WARM_RED, True),
        ("", 8, WHITE, False),
        ("Forget everything.", 22, WHITE, False),
        ("Go back to doing", 22, WHITE, False),
        ("what you've always done.", 22, WHITE, False),
        ("", 8, WHITE, False),
        ("Keep losing ₹40,000+", 22, WARM_RED, True),
        ("every month.", 22, WARM_RED, True),
    ],
    alignment=PP_ALIGN.LEFT
)

# Option 2 box
o2 = slide.shapes.add_shape(1, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5))
o2.fill.solid()
o2.fill.fore_color.rgb = RGBColor(0x0A, 0x2D, 0x12)
o2.line.color.rgb = SUCCESS_GRN
o2.line.width = Pt(2)

add_multiline_box(
    slide, Inches(7.4), Inches(2.0), Inches(4.8), Inches(4.0),
    [
        ("OPTION 2", 24, SUCCESS_GRN, True),
        ("", 8, WHITE, False),
        ("Get a proven roadmap", 22, WHITE, False),
        ("from someone who's", 22, WHITE, False),
        ("already solved this.", 22, WHITE, False),
        ("", 8, WHITE, False),
        ("Book your FREE", 22, SUCCESS_GRN, True),
        ("Store Check-Up now.", 22, SUCCESS_GRN, True),
    ],
    alignment=PP_ALIGN.LEFT
)

add_fade_in(slide, o1, delay_ms=0)
add_fade_in(slide, o2, delay_ms=600, duration_ms=1000)
add_slide_transition(slide, 'fade')

# Final CTA
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, 0x0A, 0x0E, 0x1A)

box = add_multiline_box(
    slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(2.5),
    [
        ("Click the button below this video.", 34, WHITE, True),
        ("Book your free 30-minute Store Check-Up.", 34, ACCENT_GOLD, True),
        ("", 16, WHITE, False),
        ("I'll speak to you soon.", 28, LIGHT_GRAY, False),
    ]
)

btn = slide.shapes.add_shape(
    5, Inches(3.0), Inches(4.8), Inches(7.333), Inches(1.2)
)
btn.fill.solid()
btn.fill.fore_color.rgb = SUCCESS_GRN
btn.line.fill.background()

add_text_box(slide, Inches(3.0), Inches(4.9), Inches(7.333), Inches(1.0),
    "👉  BOOK MY FREE STORE CHECK-UP  →", font_size=28, color=WHITE)

add_text_box(slide, Inches(2), Inches(6.3), Inches(9.333), Inches(0.6),
    "RETAIL CONTROL ARCHITECT™  •  Stop the Silent Bleed",
    font_size=16, color=ACCENT_GOLD, bold=False)

add_fade_in(slide, box, delay_ms=0)
add_fade_in(slide, btn, delay_ms=800, duration_ms=1000)
add_slide_transition(slide, 'fade')


# ── Save ───────────────────────────────────────────────────────────
output_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "RCA_VSL_Presentation.pptx"
)
prs.save(output_path)
print(f"✅ Successfully created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
