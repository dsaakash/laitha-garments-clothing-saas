"""
Generate a beautiful WhatsApp Cold DM Script Templates PPT.
Messages to send store owners via WhatsApp BEFORE calling/visiting.
Multiple templates for different scenarios.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette (WhatsApp-inspired + Premium Dark) ───
BG_DARK       = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD       = RGBColor(0x1A, 0x1A, 0x2E)
WA_GREEN      = RGBColor(0x25, 0xD3, 0x66)   # WhatsApp green
WA_DARK_GREEN = RGBColor(0x07, 0x5E, 0x54)   # WhatsApp dark teal
WA_LIGHT_GREEN= RGBColor(0xDC, 0xF8, 0xC6)   # WhatsApp chat bubble
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
ACCENT_GOLD   = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_GREEN  = RGBColor(0x00, 0xE6, 0x76)
ACCENT_RED    = RGBColor(0xFF, 0x4D, 0x4D)
ACCENT_BLUE   = RGBColor(0x4D, 0xA8, 0xFF)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xB0, 0xB0, 0xC0)
MEDIUM_GRAY   = RGBColor(0x80, 0x80, 0x99)
SUBTLE_PURPLE = RGBColor(0x6C, 0x5C, 0xE7)
CHAT_BG       = RGBColor(0x0B, 0x14, 0x1A)   # WhatsApp chat dark bg
CHAT_BUBBLE_SEND = RGBColor(0x00, 0x5C, 0x4B) # Sent message green
CHAT_BUBBLE_RECV = RGBColor(0x20, 0x2C, 0x33) # Received message dark

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=28,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_name="Calibri"):
    """Each line = (text, size, color, bold, spacing_after)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, spacing) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(spacing)
        p.space_before = Pt(0)
    return txBox


def add_rounded_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
    shape.line.width = Pt(1)
    return shape


def add_chat_bubble(slide, left, top, width, height, text, font_size=16,
                    is_sent=True, color=WHITE):
    """Create a WhatsApp-style chat bubble."""
    bg = CHAT_BUBBLE_SEND if is_sent else CHAT_BUBBLE_RECV
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    # Text inside bubble
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1),
                                      width - Inches(0.3), height - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return shape


def add_wa_message_card(slide, left, top, width, msg_lines, label="", label_color=WA_GREEN):
    """Full WhatsApp style message card with label and content."""
    # Calculate height
    line_count = len(msg_lines)
    total_height = Inches(0.35 * line_count + 0.6)

    # Card bg
    add_rounded_card(slide, left, top, width, total_height, CHAT_BG)

    # Label on top
    if label:
        add_text_box(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.3),
                     label, font_size=14, color=label_color, bold=True)

    # Message lines
    y = top + Inches(0.45) if label else top + Inches(0.15)
    for line in msg_lines:
        add_text_box(slide, left + Inches(0.3), y, width - Inches(0.6), Inches(0.3),
                     line, font_size=15, color=LIGHT_GRAY)
        y += Inches(0.32)

    return total_height


def add_number_badge(slide, left, top, number, color=ACCENT_ORANGE):
    size = Inches(0.65)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return circle


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, WA_GREEN)

# WhatsApp icon circle
add_circle(slide, Inches(10.2), Inches(0.8), Inches(2.3), WA_GREEN)
add_text_box(slide, Inches(10.2), Inches(1.4), Inches(2.3), Inches(1.0),
             "WA", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(1.3), Inches(9), Inches(1.2),
             "WhatsApp Cold DM", font_size=56, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.7), Inches(9), Inches(1.2),
             "Script Templates", font_size=56, color=WA_GREEN, bold=True)

add_text_box(slide, Inches(1), Inches(4.3), Inches(8), Inches(1),
             "Ready-to-use WhatsApp messages to send store owners BEFORE calling or visiting.\nJust copy, personalize the name, and send.",
             font_size=22, color=LIGHT_GRAY)

add_accent_bar(slide, Inches(1), Inches(6.0), Inches(2), Inches(0.06), WA_GREEN)
add_text_box(slide, Inches(1), Inches(6.2), Inches(6), Inches(0.5),
             "Retail Control Architect  |  Cold Outreach Kit",
             font_size=16, color=MEDIUM_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Strategy Overview — When to DM vs Call
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GOLD)

add_text_box(slide, Inches(1), Inches(0.4), Inches(10), Inches(0.8),
             "Your Outreach Strategy — Step by Step", font_size=40, color=ACCENT_GOLD, bold=True)

steps = [
    ("1", "Send WhatsApp DM First", "Send a friendly message to introduce yourself. Don't sell anything yet.", WA_GREEN, "FIRST"),
    ("2", "Wait for Reply (2-4 hours)", "If they reply — great! Start a conversation. If no reply — move to step 3.", ACCENT_BLUE, "THEN"),
    ("3", "Make the Call", "Call them and reference your WhatsApp message. They already know who you are.", ACCENT_ORANGE, "THEN"),
    ("4", "Follow Up on WhatsApp", "After the call, send a thank-you message + summary of what you discussed.", SUBTLE_PURPLE, "THEN"),
    ("5", "Book the Store Visit", "Confirm the meeting date/time on WhatsApp so it's in writing.", ACCENT_GREEN, "GOAL"),
]

for i, (num, title, desc, color, tag) in enumerate(steps):
    y = Inches(1.5) + Inches(i * 1.1)
    add_rounded_card(slide, Inches(1), y, Inches(11), Inches(0.9), BG_CARD)

    # Step number badge
    add_number_badge(slide, Inches(1.2), y + Inches(0.12), num, color)

    # Title
    add_text_box(slide, Inches(2.1), y + Inches(0.08), Inches(4), Inches(0.4),
                 title, font_size=22, color=WHITE, bold=True)

    # Description
    add_text_box(slide, Inches(2.1), y + Inches(0.48), Inches(7), Inches(0.4),
                 desc, font_size=16, color=MEDIUM_GRAY)

    # Tag
    tag_shape = add_rounded_card(slide, Inches(10.2), y + Inches(0.22), Inches(1.3), Inches(0.42), color)
    add_text_box(slide, Inches(10.2), y + Inches(0.24), Inches(1.3), Inches(0.4),
                 tag, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Template 1 — First Contact DM (Soft Approach)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, WA_GREEN)

add_text_box(slide, Inches(1), Inches(0.3), Inches(5), Inches(0.4),
             "TEMPLATE 1", font_size=16, color=WA_GREEN, bold=True)
add_text_box(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.7),
             "First Contact — Soft Introduction", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(6), Inches(0.4),
             "Use when: Reaching out for the first time", font_size=16, color=MEDIUM_GRAY)

# Chat bubble simulation
add_rounded_card(slide, Inches(1), Inches(1.9), Inches(7.5), Inches(5.0), CHAT_BG)

# Header bar (WhatsApp style)
add_accent_bar(slide, Inches(1), Inches(1.9), Inches(7.5), Inches(0.5), WA_DARK_GREEN)
add_text_box(slide, Inches(1.3), Inches(1.95), Inches(5), Inches(0.4),
             "Store Owner Name", font_size=16, color=WHITE, bold=True)

# Message bubbles
msgs = [
    ("Hi [Name] ji, Namaste!", True, Inches(2.6)),
    ("I am Aakash from Retail Control Architect.", True, Inches(3.1)),
    ("I help clothing & retail store owners find and fix hidden stock losses.", True, Inches(3.6)),
    ("I noticed your store [Store Name] and thought I could share something useful.", True, Inches(4.2)),
    ("Many store owners lose 2-5 lakhs every year from stock mismatch without even knowing.", True, Inches(4.8)),
    ("Would you be open to a quick 3-minute chat about how to check if this is happening in your store too?", True, Inches(5.4)),
    ("No pressure at all. Just thought it might help.", True, Inches(6.0)),
]

for text, is_sent, y_pos in msgs:
    x = Inches(1.3) if is_sent else Inches(4.5)
    w = Inches(6.8)
    add_chat_bubble(slide, x, y_pos, w, Inches(0.4), text, font_size=14, is_sent=is_sent)

# Tips on the right
add_rounded_card(slide, Inches(9), Inches(1.9), Inches(3.8), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(9.2), Inches(2.1), Inches(3.4), Inches(0.4),
             "Tips for this message:", font_size=18, color=ACCENT_GOLD, bold=True)

tips = [
    "Use their actual name",
    "Mention their store name",
    "Keep it friendly, not sales-y",
    "Short paragraphs = easy to read",
    "End with a simple YES/NO question",
    "Send between 10 AM - 12 PM",
    "Add a namaste - shows respect",
]

for i, tip in enumerate(tips):
    add_text_box(slide, Inches(9.3), Inches(2.7) + Inches(i * 0.5), Inches(3.3), Inches(0.5),
                 f"  {tip}", font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Template 2 — Pain Point DM (Direct Approach)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.3), Inches(5), Inches(0.4),
             "TEMPLATE 2", font_size=16, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.7),
             "Pain Point Opener — Direct Approach", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(8), Inches(0.4),
             "Use when: You want to grab attention immediately with a relatable problem", font_size=16, color=MEDIUM_GRAY)

# Chat simulation
add_rounded_card(slide, Inches(1), Inches(1.9), Inches(7.5), Inches(5.0), CHAT_BG)
add_accent_bar(slide, Inches(1), Inches(1.9), Inches(7.5), Inches(0.5), WA_DARK_GREEN)
add_text_box(slide, Inches(1.3), Inches(1.95), Inches(5), Inches(0.4),
             "Store Owner Name", font_size=16, color=WHITE, bold=True)

msgs = [
    ("Hi [Name] ji", True, Inches(2.6)),
    ("Quick question - does your store ever face this?", True, Inches(3.0)),
    ("Computer shows 100 items... but on the shelf you count only 85?", True, Inches(3.5)),
    ("I have helped 10+ store owners fix this exact problem.", True, Inches(4.1)),
    ("One owner found Rs 3.2 lakhs of hidden losses in just 45 days.", True, Inches(4.6)),
    ("I would love to show you a simple way to check if your store has the same issue.", True, Inches(5.2)),
    ("Can I call you for just 3 minutes tomorrow?", True, Inches(5.8)),
]

for text, is_sent, y_pos in msgs:
    x = Inches(1.3)
    w = Inches(6.8)
    add_chat_bubble(slide, x, y_pos, w, Inches(0.4), text, font_size=14, is_sent=is_sent)

# Why it works
add_rounded_card(slide, Inches(9), Inches(1.9), Inches(3.8), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(9.2), Inches(2.1), Inches(3.4), Inches(0.4),
             "Why this works:", font_size=18, color=ACCENT_ORANGE, bold=True)

reasons = [
    "Starts with a question they relate to",
    "Uses real numbers (100 vs 85)",
    "Shows proof (Rs 3.2 lakhs)",
    "Mentions other store owners",
    "Asks for only 3 minutes",
    "Specific time (tomorrow)",
    "No pushy language",
]

for i, reason in enumerate(reasons):
    add_text_box(slide, Inches(9.3), Inches(2.7) + Inches(i * 0.5), Inches(3.3), Inches(0.5),
                 f"  {reason}", font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Template 3 — Value-First DM (Give Before Ask)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GREEN)

add_text_box(slide, Inches(1), Inches(0.3), Inches(5), Inches(0.4),
             "TEMPLATE 3", font_size=16, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.7),
             "Value-First — Give Before You Ask", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(8), Inches(0.4),
             "Use when: You want to build trust by sharing something useful first", font_size=16, color=MEDIUM_GRAY)

# Chat simulation
add_rounded_card(slide, Inches(1), Inches(1.9), Inches(7.5), Inches(5.0), CHAT_BG)
add_accent_bar(slide, Inches(1), Inches(1.9), Inches(7.5), Inches(0.5), WA_DARK_GREEN)
add_text_box(slide, Inches(1.3), Inches(1.95), Inches(5), Inches(0.4),
             "Store Owner Name", font_size=16, color=WHITE, bold=True)

msgs = [
    ("Hi [Name] ji, I hope you are doing well!", True, Inches(2.6)),
    ("I work with retail store owners, and I wanted to share 3 quick tips that helped my clients save lakhs:", True, Inches(3.1)),
    ("1. Count your top 20 best-selling items every week (not all items, just 20)", True, Inches(3.7)),
    ("2. Make staff write down every return - even small ones", True, Inches(4.2)),
    ("3. Check closing stock every Saturday before leaving", True, Inches(4.7)),
    ("These 3 things alone can save Rs 50,000+ every month for most stores.", True, Inches(5.2)),
    ("If you want, I can do a FREE check of your store and show you exactly where money is slipping out.", True, Inches(5.8)),
    ("Just reply YES and I will arrange it.", True, Inches(6.4)),
]

for text, is_sent, y_pos in msgs:
    x = Inches(1.3)
    w = Inches(6.8)
    add_chat_bubble(slide, x, y_pos, w, Inches(0.4), text, font_size=13, is_sent=is_sent)

# Right panel
add_rounded_card(slide, Inches(9), Inches(1.9), Inches(3.8), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(9.2), Inches(2.1), Inches(3.4), Inches(0.4),
             "Power move:", font_size=18, color=ACCENT_GREEN, bold=True)

reasons = [
    "Gives FREE value upfront",
    "Shows you know their world",
    "3 actionable tips they can use",
    "Makes them think 'this guy knows'",
    "Ends with simple 'reply YES'",
    "Low commitment ask",
    "Builds trust before the call",
]

for i, reason in enumerate(reasons):
    add_text_box(slide, Inches(9.3), Inches(2.7) + Inches(i * 0.5), Inches(3.3), Inches(0.5),
                 f"  {reason}", font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Template 4 — Social Proof DM (Story-Based)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, SUBTLE_PURPLE)

add_text_box(slide, Inches(1), Inches(0.3), Inches(5), Inches(0.4),
             "TEMPLATE 4", font_size=16, color=SUBTLE_PURPLE, bold=True)
add_text_box(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.7),
             "Social Proof — Share a Real Story", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(8), Inches(0.4),
             "Use when: You want to create curiosity through a real-life example", font_size=16, color=MEDIUM_GRAY)

# Chat simulation
add_rounded_card(slide, Inches(1), Inches(1.9), Inches(7.5), Inches(5.0), CHAT_BG)
add_accent_bar(slide, Inches(1), Inches(1.9), Inches(7.5), Inches(0.5), WA_DARK_GREEN)
add_text_box(slide, Inches(1.3), Inches(1.95), Inches(5), Inches(0.4),
             "Store Owner Name", font_size=16, color=WHITE, bold=True)

msgs = [
    ("Hi [Name] ji, Namaste", True, Inches(2.6)),
    ("I wanted to share a quick story with you...", True, Inches(3.0)),
    ("Last month I visited a clothing store. The owner thought everything was fine.", True, Inches(3.5)),
    ("His computer showed stock worth Rs 28 lakhs.", True, Inches(4.0)),
    ("But when we actually counted... only Rs 24 lakhs was there.", True, Inches(4.5)),
    ("Rs 4 lakhs had silently disappeared and he had no idea!", True, Inches(5.0)),
    ("We fixed it in 30 days using a simple system.", True, Inches(5.5)),
    ("Would you like to know if something similar is happening in your store? I can check for FREE.", True, Inches(6.0)),
]

for text, is_sent, y_pos in msgs:
    x = Inches(1.3)
    w = Inches(6.8)
    add_chat_bubble(slide, x, y_pos, w, Inches(0.4), text, font_size=14, is_sent=is_sent)

# Right panel
add_rounded_card(slide, Inches(9), Inches(1.9), Inches(3.8), Inches(5.0), BG_CARD)
add_text_box(slide, Inches(9.2), Inches(2.1), Inches(3.4), Inches(0.4),
             "Story selling works:", font_size=18, color=SUBTLE_PURPLE, bold=True)

reasons = [
    "Real story creates curiosity",
    "Specific numbers feel authentic",
    "Rs 28L vs Rs 24L = dramatic gap",
    "Owner can relate to the situation",
    "'Silently disappeared' = scary",
    "'Fixed in 30 days' = hope",
    "FREE check = no risk to say yes",
]

for i, reason in enumerate(reasons):
    add_text_box(slide, Inches(9.3), Inches(2.7) + Inches(i * 0.5), Inches(3.3), Inches(0.5),
                 f"  {reason}", font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Template 5 — Follow-Up After No Reply
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(0.3), Inches(5), Inches(0.4),
             "TEMPLATE 5", font_size=16, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.7),
             "Follow-Up — When They Did Not Reply", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(8), Inches(0.4),
             "Use when: 2-3 days have passed since your first message with no reply", font_size=16, color=MEDIUM_GRAY)

# Two follow-up options side by side
# Option A
add_rounded_card(slide, Inches(0.8), Inches(2.0), Inches(5.7), Inches(5.0), CHAT_BG)
add_accent_bar(slide, Inches(0.8), Inches(2.0), Inches(5.7), Inches(0.45), WA_DARK_GREEN)
add_text_box(slide, Inches(1.0), Inches(2.05), Inches(3), Inches(0.35),
             "Follow-Up Option A (Gentle)", font_size=14, color=WHITE, bold=True)

followup_a = [
    ("Hi [Name] ji", Inches(2.7)),
    ("Just checking if you got my last message.", Inches(3.1)),
    ("No worries if you are busy!", Inches(3.5)),
    ("I just wanted to share that I am offering a FREE store stock check this week.", Inches(3.9)),
    ("It takes only 30 minutes and can save you lakhs.", Inches(4.4)),
    ("Let me know if you are interested.", Inches(4.8)),
    ("Have a great day!", Inches(5.2)),
]

for text, y_pos in followup_a:
    add_chat_bubble(slide, Inches(1.0), y_pos, Inches(5.2), Inches(0.35),
                    text, font_size=13, is_sent=True)

# Option B
add_rounded_card(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(5.0), CHAT_BG)
add_accent_bar(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(0.45), ACCENT_ORANGE)
add_text_box(slide, Inches(7.0), Inches(2.05), Inches(3.5), Inches(0.35),
             "Follow-Up Option B (Curiosity)", font_size=14, color=WHITE, bold=True)

followup_b = [
    ("Hi [Name] ji", Inches(2.7)),
    ("I visited 3 stores in your area last week.", Inches(3.1)),
    ("Every single one had the same problem -", Inches(3.5)),
    ("stock on computer and stock on shelf did not match.", Inches(3.9)),
    ("One owner was losing Rs 40,000 EVERY MONTH without knowing.", Inches(4.4)),
    ("I can check if this is happening in your store too. 100% free.", Inches(4.9)),
    ("Interested?", Inches(5.3)),
]

for text, y_pos in followup_b:
    add_chat_bubble(slide, Inches(7.0), y_pos, Inches(5.2), Inches(0.35),
                    text, font_size=13, is_sent=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Template 6 — After Phone Call Thank You
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GOLD)

add_text_box(slide, Inches(1), Inches(0.3), Inches(5), Inches(0.4),
             "TEMPLATE 6", font_size=16, color=ACCENT_GOLD, bold=True)
add_text_box(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.7),
             "After The Call — Thank You + Confirm Visit", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(8), Inches(0.4),
             "Use when: Right after a successful phone call to lock in the meeting", font_size=16, color=MEDIUM_GRAY)

# Chat simulation
add_rounded_card(slide, Inches(1), Inches(1.9), Inches(7.5), Inches(5.2), CHAT_BG)
add_accent_bar(slide, Inches(1), Inches(1.9), Inches(7.5), Inches(0.5), WA_DARK_GREEN)
add_text_box(slide, Inches(1.3), Inches(1.95), Inches(5), Inches(0.4),
             "Store Owner Name", font_size=16, color=WHITE, bold=True)

msgs = [
    ("Thank you so much for your time [Name] ji! Really appreciate it.", True, Inches(2.6)),
    ("As we discussed, here is what I will do when I visit your store:", True, Inches(3.1)),
    ("1. Quick walk-through of your stock area", True, Inches(3.6)),
    ("2. Check a few items: computer count vs actual count", True, Inches(4.0)),
    ("3. Show you where the gaps are", True, Inches(4.4)),
    ("4. Give you a clear picture + action plan", True, Inches(4.8)),
    ("All FREE - no charges, no obligation.", True, Inches(5.3)),
    ("Confirmed: [Day], [Date] at [Time]", True, Inches(5.8)),
    ("Looking forward to meeting you!", True, Inches(6.3)),
]

for text, is_sent, y_pos in msgs:
    add_chat_bubble(slide, Inches(1.3), y_pos, Inches(6.8), Inches(0.35),
                    text, font_size=13, is_sent=is_sent)

# Right tips
add_rounded_card(slide, Inches(9), Inches(1.9), Inches(3.8), Inches(5.2), BG_CARD)
add_text_box(slide, Inches(9.2), Inches(2.1), Inches(3.4), Inches(0.4),
             "Why send this:", font_size=18, color=ACCENT_GOLD, bold=True)

reasons = [
    "Shows professionalism",
    "Confirms the meeting in writing",
    "Reminds them what to expect",
    "'FREE' removes any last doubt",
    "They can share with partner/staff",
    "Creates psychological commitment",
    "Hard to cancel after this!",
]

for i, reason in enumerate(reasons):
    add_text_box(slide, Inches(9.3), Inches(2.7) + Inches(i * 0.5), Inches(3.3), Inches(0.5),
                 f"  {reason}", font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Template 7 — Meeting Reminder (Day Before)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, WA_GREEN)

add_text_box(slide, Inches(1), Inches(0.3), Inches(5), Inches(0.4),
             "TEMPLATE 7", font_size=16, color=WA_GREEN, bold=True)
add_text_box(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.7),
             "Day-Before Reminder + After Visit Follow-Up", font_size=36, color=WHITE, bold=True)

# Two templates side by side
# Reminder
add_rounded_card(slide, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0), CHAT_BG)
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(5.7), Inches(0.45), WA_DARK_GREEN)
add_text_box(slide, Inches(1.0), Inches(1.85), Inches(4), Inches(0.35),
             "Day-Before Reminder", font_size=14, color=WHITE, bold=True)

reminder_msgs = [
    ("Hi [Name] ji, good evening!", Inches(2.5)),
    ("Just a reminder about our meeting tomorrow", Inches(2.9)),
    ("at [Time] at your store.", Inches(3.3)),
    ("I will bring everything needed for the stock check.", Inches(3.7)),
    ("You do not need to prepare anything.", Inches(4.1)),
    ("If the timing needs to change, just let me know!", Inches(4.5)),
    ("See you tomorrow!", Inches(4.9)),
]

for text, y_pos in reminder_msgs:
    add_chat_bubble(slide, Inches(1.0), y_pos, Inches(5.2), Inches(0.35),
                    text, font_size=13, is_sent=True)

# After visit
add_rounded_card(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0), CHAT_BG)
add_accent_bar(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(0.45), ACCENT_GREEN)
add_text_box(slide, Inches(7.0), Inches(1.85), Inches(4), Inches(0.35),
             "After Visit Follow-Up", font_size=14, color=WHITE, bold=True)

after_msgs = [
    ("[Name] ji, thank you for today!", Inches(2.5)),
    ("It was great meeting you and seeing your store.", Inches(2.9)),
    ("As I showed you, here are the key findings:", Inches(3.3)),
    ("- [Finding 1: e.g., 15% stock mismatch found]", Inches(3.7)),
    ("- [Finding 2: e.g., Returns not being recorded]", Inches(4.1)),
    ("- [Finding 3: e.g., No weekly count process]", Inches(4.5)),
    ("I am confident we can fix this in 30 days.", Inches(4.9)),
    ("Shall we start next week?", Inches(5.3)),
]

for text, y_pos in after_msgs:
    add_chat_bubble(slide, Inches(7.0), y_pos, Inches(5.2), Inches(0.35),
                    text, font_size=13, is_sent=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Cold Call Phone Script — What to Say
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.3), Inches(5), Inches(0.4),
             "CALL SCRIPT", font_size=16, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(1), Inches(0.7), Inches(10), Inches(0.7),
             "Phone Call Script — After the WhatsApp DM", font_size=36, color=WHITE, bold=True)

# Script in sections
sections = [
    ("OPENING (15 seconds)", ACCENT_BLUE, [
        "Hi, am I speaking with [Name] ji?",
        "This is Aakash from Retail Control Architect.",
        "I had sent you a WhatsApp message yesterday about stock management.",
        "Do you have 3 minutes? I promise I will be quick.",
    ]),
    ("PROBLEM (30 seconds)", ACCENT_RED, [
        "I work with clothing store owners who face a common problem -",
        "their computer shows one stock number, but the real count is different.",
        "Most owners think this is normal... but it actually means money is leaking silently.",
        "One store I worked with was losing Rs 4 lakhs and did not even know!",
    ]),
    ("SOLUTION (30 seconds)", ACCENT_GREEN, [
        "I have a simple 5-step system that fixes this in 30 days.",
        "No expensive software needed. Just a clear process for your store.",
        "I would love to visit your store for just 30 minutes - completely free -",
        "and show you exactly where your store might be losing money.",
    ]),
    ("CLOSE (15 seconds)", ACCENT_GOLD, [
        "Would [Tuesday/Thursday] work for a quick 30-minute visit?",
        "There is absolutely no charge and no obligation.",
        "I just want to help you see what most store owners are missing.",
    ]),
]

y = Inches(1.5)
for section_title, color, lines in sections:
    add_rounded_card(slide, Inches(1), y, Inches(11), Inches(0.3 + len(lines) * 0.35), BG_CARD)
    add_accent_bar(slide, Inches(1), y, Inches(0.1), Inches(0.3 + len(lines) * 0.35), color)
    add_text_box(slide, Inches(1.3), y + Inches(0.05), Inches(3), Inches(0.3),
                 section_title, font_size=14, color=color, bold=True)
    for i, line in enumerate(lines):
        add_text_box(slide, Inches(1.5), y + Inches(0.35 + i * 0.3), Inches(10), Inches(0.3),
                     f'"{line}"', font_size=15, color=LIGHT_GRAY)
    y += Inches(0.5 + len(lines) * 0.33)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: Handling Replies — What If They Say...
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, SUBTLE_PURPLE)

add_text_box(slide, Inches(1), Inches(0.3), Inches(10), Inches(0.7),
             "What If They Reply With... (Reply Templates)", font_size=36, color=WHITE, bold=True)

replies = [
    ("They say: 'Who are you?'",
     "I am Aakash. I help store owners find hidden stock losses. I am not selling any software or product. I just want to show you something useful for your store.",
     ACCENT_BLUE),
    ("They say: 'I am busy'",
     "No problem at all [Name] ji! When would be a better time? I need just 3 minutes. Even a quick call on [Day] evening would work.",
     ACCENT_ORANGE),
    ("They say: 'How much does it cost?'",
     "The store check-up is 100% free. If you like what I show you, we can discuss options later. But the first step costs you nothing.",
     ACCENT_GREEN),
    ("They say: 'I already use software'",
     "That is great! But software alone does not stop stock leakage. My system adds the missing piece - a PROCESS that makes sure everything is tracked properly. Let me show you what I mean.",
     ACCENT_GOLD),
    ("They say: 'Not interested'",
     "I completely understand [Name] ji. Just one last thing - if you ever notice that your stock count does not match, feel free to message me. I will be happy to help. Thank you for your time!",
     ACCENT_RED),
    ("They say: 'Send me more details'",
     "Sure! Let me share a quick summary. But honestly, the best way to understand is a 30-minute visit where I show you live with your own stock. It is much more powerful than reading. Can we schedule that?",
     SUBTLE_PURPLE),
]

for i, (their_reply, your_response, color) in enumerate(replies):
    y = Inches(1.2) + Inches(i * 1.0)
    add_rounded_card(slide, Inches(0.8), y, Inches(11.5), Inches(0.85), BG_CARD)
    add_accent_bar(slide, Inches(0.8), y, Inches(0.08), Inches(0.85), color)
    add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(3.5), Inches(0.35),
                 their_reply, font_size=15, color=color, bold=True)
    add_text_box(slide, Inches(4.5), y + Inches(0.05), Inches(7.5), Inches(0.75),
                 your_response, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: Daily Outreach Plan
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GREEN)

add_text_box(slide, Inches(1), Inches(0.3), Inches(10), Inches(0.7),
             "Your Daily Outreach Plan", font_size=38, color=ACCENT_GREEN, bold=True)

# Time blocks
time_blocks = [
    ("9:00 AM - 10:00 AM", "Send 10-15 WhatsApp DMs",
     "Use Templates 1-4. Personalize each message with their name and store name.",
     WA_GREEN),
    ("10:00 AM - 10:30 AM", "Follow up on yesterday's DMs",
     "Check who read but did not reply. Send Template 5 follow-up messages.",
     ACCENT_BLUE),
    ("10:30 AM - 12:00 PM", "Make phone calls",
     "Call everyone who replied positively. Use the phone script from Slide 10.",
     ACCENT_ORANGE),
    ("2:00 PM - 3:00 PM", "Send after-call messages",
     "Send Template 6 to everyone you spoke with. Confirm store visits.",
     ACCENT_GOLD),
    ("5:00 PM - 5:30 PM", "Send reminders",
     "Template 7 reminders for tomorrow's scheduled store visits.",
     SUBTLE_PURPLE),
    ("Evening", "Update your tracker",
     "Log all conversations, scores (Green/Yellow/Red), and next steps.",
     ACCENT_GREEN),
]

for i, (time, activity, detail, color) in enumerate(time_blocks):
    y = Inches(1.2) + Inches(i * 1.0)
    add_rounded_card(slide, Inches(0.8), y, Inches(11.5), Inches(0.85), BG_CARD)
    add_accent_bar(slide, Inches(0.8), y, Inches(0.08), Inches(0.85), color)
    add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(2.5), Inches(0.35),
                 time, font_size=15, color=color, bold=True)
    add_text_box(slide, Inches(3.7), y + Inches(0.05), Inches(3), Inches(0.35),
                 activity, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.7), y + Inches(0.42), Inches(8), Inches(0.35),
                 detail, font_size=13, color=MEDIUM_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13: Do's & Don'ts for WhatsApp
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GOLD)

add_text_box(slide, Inches(1), Inches(0.3), Inches(10), Inches(0.7),
             "WhatsApp Do's and Don'ts", font_size=38, color=ACCENT_GOLD, bold=True)

# DOs
add_rounded_card(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), RGBColor(0x0A, 0x2A, 0x15))
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(5), Inches(0.5),
             "DO's", font_size=28, color=ACCENT_GREEN, bold=True)

dos = [
    "Use their name - makes it personal",
    "Keep messages short (under 10 lines)",
    "Use simple Hindi-English mix if needed",
    "Send during business hours only",
    "Wait 2-3 days before follow-up",
    "Always be polite and respectful",
    "Use emojis sparingly (1-2 max)",
    "End with a clear question",
    "Save their contact properly",
    "Track who replied, who didn't",
]

for i, text in enumerate(dos):
    add_text_box(slide, Inches(1.2), Inches(2.2) + Inches(i * 0.44), Inches(4.8), Inches(0.4),
                 f"  {text}", font_size=14, color=LIGHT_GRAY)

# DON'Ts
add_rounded_card(slide, Inches(6.8), Inches(1.3), Inches(5.5), Inches(5.5), RGBColor(0x2A, 0x15, 0x15))
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5), Inches(0.5),
             "DON'Ts", font_size=28, color=ACCENT_RED, bold=True)

donts = [
    "Never send the same message to everyone",
    "Never send voice notes on first contact",
    "Never share price in the first message",
    "Never send long paragraphs (wall of text)",
    "Never call without messaging first",
    "Never spam with multiple messages same day",
    "Never use technical words they won't understand",
    "Never argue if they say not interested",
    "Never send messages late at night",
    "Never add them to groups without asking",
]

for i, text in enumerate(donts):
    add_text_box(slide, Inches(7.2), Inches(2.2) + Inches(i * 0.44), Inches(4.8), Inches(0.4),
                 f"  {text}", font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14: Quick Copy-Paste Messages (All in One)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, WA_GREEN)

add_text_box(slide, Inches(1), Inches(0.3), Inches(10), Inches(0.7),
             "Quick-Copy Messages (Short Versions)", font_size=38, color=WA_GREEN, bold=True)
add_text_box(slide, Inches(1), Inches(0.9), Inches(10), Inches(0.4),
             "Shorter versions you can quickly copy-paste and personalize", font_size=16, color=MEDIUM_GRAY)

# Short templates
short_templates = [
    ("First Contact (Short)", WA_GREEN,
     "Hi [Name] ji. I am Aakash. I help store owners find hidden stock losses. Many stores lose 2-5 lakhs yearly without knowing. Can I share how to check this in your store? Just 3 min."),
    ("Pain Point (Short)", ACCENT_ORANGE,
     "Hi [Name] ji. Does your stock count ever not match the computer? I helped a store owner recover Rs 3.2 lakhs in 45 days. Would you like to know how? I can explain in 3 minutes."),
    ("Value-First (Short)", ACCENT_GREEN,
     "Hi [Name] ji. Quick tip: Count your top 20 items weekly + make staff record every return. This alone saves Rs 50K/month. Want me to check your store for FREE? Reply YES."),
    ("Follow-Up (Short)", ACCENT_BLUE,
     "Hi [Name] ji, just following up. I am doing FREE store stock check-ups this week. Takes 30 min, can save lakhs. Want me to visit your store? No cost, no pressure."),
    ("Confirm Visit (Short)", ACCENT_GOLD,
     "Thank you [Name] ji! Confirmed: [Day, Date, Time] at your store. I will check your stock, find the gaps, and give you a clear plan. All free. See you soon!"),
    ("Day-Before Reminder (Short)", SUBTLE_PURPLE,
     "Hi [Name] ji! Just a reminder - I will be at your store tomorrow at [Time]. You do not need to prepare anything. See you!"),
]

for i, (label, color, msg) in enumerate(short_templates):
    y = Inches(1.5) + Inches(i * 0.95)
    add_rounded_card(slide, Inches(0.8), y, Inches(11.5), Inches(0.8), BG_CARD)
    add_accent_bar(slide, Inches(0.8), y, Inches(0.08), Inches(0.8), color)
    add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(2.8), Inches(0.3),
                 label, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(1.1), y + Inches(0.33), Inches(10.8), Inches(0.45),
                 msg, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15: Summary — Your Complete Outreach Kit
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.3), Inches(10), Inches(0.7),
             "Your Complete Outreach Kit - Summary", font_size=38, color=ACCENT_ORANGE, bold=True)

kit_items = [
    ("1", "WhatsApp DM Templates (7 Templates)", "First contact, pain point, value-first, story, follow-up, after-call, reminder", WA_GREEN),
    ("2", "Phone Call Script", "90-second structured script: Opening - Problem - Solution - Close", ACCENT_ORANGE),
    ("3", "Reply Handlers", "Ready responses for 6 common replies (busy, cost, not interested, etc.)", ACCENT_BLUE),
    ("4", "Daily Outreach Plan", "Exact time blocks for DMs, calls, follow-ups, and tracking", ACCENT_GOLD),
    ("5", "Quick-Copy Messages", "Short versions of all templates for fast sending", SUBTLE_PURPLE),
]

for i, (num, title, desc, color) in enumerate(kit_items):
    y = Inches(1.3) + Inches(i * 1.0)
    add_rounded_card(slide, Inches(1), y, Inches(11), Inches(0.85), BG_CARD)
    add_number_badge(slide, Inches(1.2), y + Inches(0.1), num, color)
    add_text_box(slide, Inches(2.1), y + Inches(0.08), Inches(6), Inches(0.35),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.1), y + Inches(0.45), Inches(9), Inches(0.35),
                 desc, font_size=15, color=MEDIUM_GRAY)

# Bottom motivational
add_rounded_card(slide, Inches(2.5), Inches(6.3), Inches(8), Inches(0.8), ACCENT_ORANGE)
add_text_box(slide, Inches(2.7), Inches(6.35), Inches(7.5), Inches(0.35),
             "Target: 10-15 new DMs every day = 3-5 store visits per week", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.7), Inches(6.7), Inches(7.5), Inches(0.3),
             "Consistency wins. Keep going!", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)


# ─── Save ───
output_path = "/Users/aakash/Desktop/Week_3/WhatsApp_Cold_DM_Scripts.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
