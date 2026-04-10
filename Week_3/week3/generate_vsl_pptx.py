#!/usr/bin/env python3
"""
Generate a polished VSL (Video Sales Letter) PPTX presentation
for client demo — Retail Control Architect™
Dark theme, one-line-per-slide VSL style, 16:9 widescreen
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml
import os

# ── Constants ───────────────────────────────────────────────────
SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

# Color palette - Dark premium theme
BG_DARK      = RGBColor(0x0A, 0x0E, 0x17)   # near-black navy
BG_DARK2     = RGBColor(0x0F, 0x17, 0x27)   # slightly lighter
GOLD         = RGBColor(0xF0, 0xBD, 0x4E)   # warm gold
TEAL         = RGBColor(0x2E, 0xCC, 0xAD)   # bright teal
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xA0, 0xA8, 0xB8)
RED_WARM     = RGBColor(0xFF, 0x5A, 0x5A)
GREEN_OK     = RGBColor(0x4E, 0xD8, 0x7B)
BLUE_ACC     = RGBColor(0x5B, 0x9B, 0xF7)
ORANGE       = RGBColor(0xFF, 0x9F, 0x43)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper Functions ────────────────────────────────────────────

def set_slide_bg(slide, r, g, b):
    """Set solid background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_gradient_bg(slide, color1_hex="0A0E17", color2_hex="121E33"):
    """Set gradient background using XML."""
    bg = slide.background
    bgPr = bg._element
    # Remove existing background
    for child in list(bgPr):
        bgPr.remove(child)
    # Add solid fill as fallback (python-pptx limitation on gradients)
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(
        int(color1_hex[:2], 16),
        int(color1_hex[2:4], 16),
        int(color1_hex[4:6], 16)
    )

def add_text_slide(text, font_size=44, color=WHITE, bold=True,
                   bg_color=None, subtitle=None, sub_color=None,
                   sub_size=24, align=PP_ALIGN.CENTER,
                   accent_line=False, accent_color=None):
    """Create a single-text centered slide (VSL style)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    if bg_color:
        set_slide_bg(slide, bg_color[0], bg_color[1], bg_color[2])
    else:
        set_slide_bg(slide, 0x0A, 0x0E, 0x17)

    # Accent line at top
    if accent_line and accent_color:
        line_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), SLIDE_W, Inches(0.06)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = accent_color
        line_shape.line.fill.background()

    # Main text
    txBox = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(1.5) if subtitle else Inches(2.0),
        Inches(11.333),
        Inches(3.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = 'Calibri'

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            Inches(1.5),
            Inches(4.8),
            Inches(10.333),
            Inches(2.0)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(sub_size)
        run2.font.color.rgb = sub_color or LIGHT_GRAY
        run2.font.bold = False
        run2.font.name = 'Calibri'

    return slide

def add_split_slide(left_text, right_items, title=None,
                    left_color=WHITE, right_color=LIGHT_GRAY):
    """Create a two-column slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0x0A, 0x0E, 0x17)

    # Top accent
    line_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    line_shape.fill.solid(); line_shape.fill.fore_color.rgb = GOLD
    line_shape.line.fill.background()

    if title:
        txT = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.0))
        tf = txT.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.size = Pt(32); r.font.color.rgb = GOLD
        r.font.bold = True; r.font.name = 'Calibri'

    # Left side
    y_start = 1.8 if title else 1.2
    txL = slide.shapes.add_textbox(
        Inches(0.8), Inches(y_start), Inches(5.5), Inches(5.0)
    )
    tfL = txL.text_frame; tfL.word_wrap = True
    pL = tfL.paragraphs[0]; pL.alignment = PP_ALIGN.LEFT
    rL = pL.add_run(); rL.text = left_text; rL.font.size = Pt(26)
    rL.font.color.rgb = left_color; rL.font.bold = True; rL.font.name = 'Calibri'

    # Right side - bullet items
    txR = slide.shapes.add_textbox(
        Inches(7.0), Inches(y_start), Inches(5.5), Inches(5.0)
    )
    tfR = txR.text_frame; tfR.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tfR.paragraphs[0]
        else:
            p = tfR.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
        r = p.add_run(); r.text = item; r.font.size = Pt(20)
        r.font.color.rgb = right_color; r.font.name = 'Calibri'

    return slide

def add_number_slide(number, label, sub_label=None,
                     num_color=GOLD, label_color=WHITE):
    """Big number impact slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0x0A, 0x0E, 0x17)

    # Number
    txN = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(3.0))
    tf = txN.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = number; r.font.size = Pt(96); r.font.color.rgb = num_color
    r.font.bold = True; r.font.name = 'Calibri'

    # Label
    txL = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.333), Inches(1.5))
    tfL = txL.text_frame; pL = tfL.paragraphs[0]; pL.alignment = PP_ALIGN.CENTER
    rL = pL.add_run(); rL.text = label; rL.font.size = Pt(32); rL.font.color.rgb = label_color
    rL.font.bold = True; rL.font.name = 'Calibri'

    if sub_label:
        txS = slide.shapes.add_textbox(Inches(2.0), Inches(5.5), Inches(9.333), Inches(1.2))
        tfS = txS.text_frame; pS = tfS.paragraphs[0]; pS.alignment = PP_ALIGN.CENTER
        rS = pS.add_run(); rS.text = sub_label; rS.font.size = Pt(20)
        rS.font.color.rgb = LIGHT_GRAY; rS.font.name = 'Calibri'

    return slide

def add_step_slide(step_num, step_name, description, icon_color=TEAL):
    """Framework step slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, 0x0A, 0x0E, 0x17)

    # Top accent
    line = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = icon_color; line.line.fill.background()

    # Step circle (simulated with text)
    txC = slide.shapes.add_textbox(Inches(5.166), Inches(0.8), Inches(3.0), Inches(1.5))
    tfC = txC.text_frame; pC = tfC.paragraphs[0]; pC.alignment = PP_ALIGN.CENTER
    rC = pC.add_run(); rC.text = f"STEP {step_num}"; rC.font.size = Pt(20)
    rC.font.color.rgb = icon_color; rC.font.bold = True; rC.font.name = 'Calibri'

    # Step name
    txN = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.333), Inches(1.5))
    tfN = txN.text_frame; pN = tfN.paragraphs[0]; pN.alignment = PP_ALIGN.CENTER
    rN = pN.add_run(); rN.text = step_name; rN.font.size = Pt(40)
    rN.font.color.rgb = WHITE; rN.font.bold = True; rN.font.name = 'Calibri'

    # Description
    txD = slide.shapes.add_textbox(Inches(2.0), Inches(4.0), Inches(9.333), Inches(2.5))
    tfD = txD.text_frame; tfD.word_wrap = True
    pD = tfD.paragraphs[0]; pD.alignment = PP_ALIGN.CENTER
    rD = pD.add_run(); rD.text = description; rD.font.size = Pt(22)
    rD.font.color.rgb = LIGHT_GRAY; rD.font.name = 'Calibri'

    return slide


# ════════════════════════════════════════════════════════════════
#                        VSL SLIDES
# ════════════════════════════════════════════════════════════════

# ── SECTION 1: THE HOOK (Slides 1–16) ──────────────────────────

add_text_slide(
    "Let me guess…",
    font_size=48, color=WHITE, bold=True
)

add_text_slide(
    "You think stock mismatch\nis normal?",
    font_size=50, color=GOLD, bold=True,
    subtitle='"Retail mein thoda bahut toh chalta hai…"',
    sub_color=LIGHT_GRAY, sub_size=28
)

add_text_slide(
    "🪣",
    font_size=96, color=WHITE, bold=False,
    subtitle="Your store is a leaking bucket.",
    sub_color=RED_WARM, sub_size=36
)

add_text_slide(
    "Every drop…\nis your profit.",
    font_size=52, color=RED_WARM, bold=True
)

add_text_slide(
    "In the next few minutes…",
    font_size=44, color=WHITE,
    subtitle="I'll show you how to stop this leak\n— in just 30 days.",
    sub_color=TEAL, sub_size=32
)

add_text_slide(
    "Without buying expensive ERP.\nWithout hiring more staff.\nWithout changing your billing software.",
    font_size=32, color=LIGHT_GRAY, bold=False,
    accent_line=True, accent_color=GOLD
)

# Quick questions - pain agitation
add_text_slide(
    "Quick question…",
    font_size=48, color=GOLD
)

add_text_slide(
    "Does your physical stock\never mismatch system stock?",
    font_size=42, color=WHITE,
    subtitle="End of month shock?  •  Staff confusion?  •  Random shortages?",
    sub_color=LIGHT_GRAY, sub_size=24
)

add_text_slide(
    "That stress…\nnever leaves you.",
    font_size=48, color=RED_WARM
)

add_text_slide(
    "You don't know what's sold…\nor what's stolen.",
    font_size=40, color=WHITE,
    subtitle="And the worst part?\nYou've accepted it as \"normal.\"",
    sub_color=LIGHT_GRAY, sub_size=26
)

add_text_slide(
    "If this sounds familiar…",
    font_size=44, color=WHITE,
    subtitle="Then what I'm about to show you\nwill change your store forever.",
    sub_color=TEAL, sub_size=30
)

# Qualifier
add_text_slide(
    "But listen carefully.",
    font_size=48, color=GOLD,
    subtitle="This is NOT for careless store owners.\nI only work with serious retail operators.",
    sub_color=LIGHT_GRAY, sub_size=24
)


# ── SECTION 2: PROOF & CREDIBILITY (Slides 17–28) ──────────────

add_text_slide(
    "Before I explain how…\nlet me show you results.",
    font_size=40, color=WHITE,
    accent_line=True, accent_color=TEAL
)

add_number_slide(
    "₹3.2L", "Recovered in 45 Days",
    "From just ONE clothing store. Hidden leakage that was invisible for years.",
    num_color=GREEN_OK
)

add_number_slide(
    "17% → 2%", "Stock Mismatch Reduced",
    "From chaos to near-perfect accuracy. In 30 days. With full compliance.",
    num_color=TEAL
)

add_number_slide(
    "0", "Stock Shocks at Month-End",
    "No more unpleasant surprises. No more guesswork. Just clean numbers.",
    num_color=GREEN_OK
)

# Who am I
add_text_slide(
    "Hi, I'm Aakash Savant.",
    font_size=48, color=WHITE,
    subtitle="Founder of Retail Control Architect™\nI help clothing stores plug hidden leakage.",
    sub_color=GOLD, sub_size=26,
    accent_line=True, accent_color=GOLD
)

add_text_slide(
    "I specialize in one thing:",
    font_size=40, color=WHITE,
    subtitle="Making your system stock\nmatch your real stock.",
    sub_color=TEAL, sub_size=36
)

# Story
add_text_slide(
    "I once walked into a store…",
    font_size=44, color=WHITE,
    subtitle="A busy clothing shop. Good revenue. Normal-looking operations.",
    sub_color=LIGHT_GRAY, sub_size=24
)

add_text_slide(
    "System showed ₹28 Lakh inventory.\nPhysical count showed ₹24 Lakh.",
    font_size=36, color=WHITE
)

add_number_slide(
    "₹4 Lakh", "Vanished.",
    'And the owner thought it was "normal."',
    num_color=RED_WARM
)

add_text_slide(
    "That's when I understood\nthe real problem.",
    font_size=44, color=GOLD
)


# ── SECTION 3: THE REAL PROBLEM (Slides 29–38) ─────────────────

add_text_slide(
    "It's not software.",
    font_size=52, color=WHITE
)

add_text_slide(
    "It's not staff.",
    font_size=52, color=WHITE
)

add_text_slide(
    "It's lack of STRUCTURE.",
    font_size=52, color=RED_WARM, bold=True
)

add_text_slide(
    "Retail chaos is the enemy.",
    font_size=48, color=GOLD,
    subtitle="Uncontrolled stock movement destroys profit silently.",
    sub_color=LIGHT_GRAY, sub_size=24
)

add_split_slide(
    "Where Does\nStock Leak?",
    [
        "❌  Sales without proper entry",
        "❌  Returns without tagging",
        "❌  Supplier inward without verification",
        "❌  Manual register alongside billing",
        "❌  End-of-day batch entries",
        "❌  Staff 'adjustments' without record",
    ],
    title="THE 6 SILENT LEAKAGE POINTS"
)

add_text_slide(
    "Small gaps\ncreate BIG losses.",
    font_size=48, color=RED_WARM,
    subtitle="₹500/day × 365 days = ₹1,82,500/year\nFrom just ONE leakage point.",
    sub_color=LIGHT_GRAY, sub_size=24
)


# ── SECTION 4: THE SOLUTION - FRAMEWORK (Slides 39–55) ─────────

add_text_slide(
    "Introducing",
    font_size=36, color=LIGHT_GRAY,
    subtitle="The 30-Day Stock Certainty System™",
    sub_color=GOLD, sub_size=48
)

add_text_slide(
    "A 7-Step Control Installation\nthat fixes stock mismatch permanently.",
    font_size=36, color=WHITE,
    accent_line=True, accent_color=TEAL,
    subtitle="Not software. Not ERP. A structural installation.",
    sub_color=LIGHT_GRAY, sub_size=24
)

# 7 Steps
add_step_slide("1", "Control Gap Audit™",
    "We measure your REAL mismatch.\nSystem stock vs physical stock.\nYour exact leakage — in Rupees.",
    TEAL)

add_step_slide("2", "Inventory Foundation Reset",
    "We clean your item structure.\nSKU naming, categories, sizes, colors.\nRemove duplicates. Fix the foundation.",
    BLUE_ACC)

add_step_slide("3", "Supplier Entry Lock™",
    "No stock enters your store\nwithout entering the system FIRST.\nYour digital gatekeeper.",
    ORANGE)

add_step_slide("4", "Sales Deduction Lock™",
    "Every sale automatically reduces stock.\nNo manual adjustments. No batch entries.\nReal-time accuracy.",
    GREEN_OK)

add_step_slide("5", "Single System Enforcement",
    "We physically REMOVE manual registers.\nNo parallel systems. No Excel.\nOne system. One truth.",
    RED_WARM)

add_step_slide("6", "Owner Visibility Dashboard",
    "YOU can check everything independently.\nStock. Suppliers. Sales. Daily report.\n10 minutes. Full control.",
    GOLD)

add_step_slide("7", "30-Day Discipline Installation",
    "We monitor. We correct. We enforce.\nDaily compliance checks for 30 days.\nUntil discipline becomes permanent habit.",
    TEAL)

# Framework summary
add_text_slide(
    "We fix behavior…\nnot just numbers.",
    font_size=44, color=GOLD,
    subtitle="You stop guessing.\nYou start controlling.",
    sub_color=WHITE, sub_size=30
)


# ── SECTION 5: RESULTS (Slides 56–62) ──────────────────────────

add_text_slide(
    "Real Results.\nReal Stores.",
    font_size=48, color=WHITE,
    accent_line=True, accent_color=GREEN_OK
)

add_number_slide(
    "98%", "Stock Accuracy Achieved",
    "From 83% to 98% in just 30 days of disciplined enforcement.",
    num_color=GREEN_OK
)

add_text_slide(
    "Internal shrinkage — eliminated.\nSecond branch — opened confidently.\nMicromanaging — stopped.",
    font_size=32, color=WHITE,
    subtitle="All without replacing software or firing staff.",
    sub_color=TEAL, sub_size=24
)


# ── SECTION 6: FUTURE VISION (Slides 63–68) ────────────────────

add_text_slide(
    "Now imagine YOUR store…",
    font_size=48, color=GOLD
)

add_text_slide(
    "Closing the month without fear.",
    font_size=44, color=WHITE,
    subtitle="Just clean reports. Accurate numbers. Zero stress.",
    sub_color=GREEN_OK, sub_size=26
)

add_text_slide(
    "Trusting your numbers.",
    font_size=44, color=WHITE,
    subtitle="Without calling your manager.\nWithout double-checking manually.",
    sub_color=GREEN_OK, sub_size=26
)

add_text_slide(
    "Sleeping peacefully.",
    font_size=48, color=TEAL,
    subtitle="Because your stock matches your shop.\nEvery single day.",
    sub_color=WHITE, sub_size=26
)


# ── SECTION 7: OBJECTION HANDLING (Slides 69–80) ───────────────

add_text_slide(
    "You might be thinking…",
    font_size=44, color=WHITE
)

# Objection 1
add_text_slide(
    '"My staff won\'t follow the system."',
    font_size=38, color=LIGHT_GRAY, bold=False
)
add_text_slide(
    "That's exactly why\nenforcement is built-in.",
    font_size=42, color=TEAL,
    subtitle="We don't just teach. We monitor and enforce\nfor 30 days straight.",
    sub_color=LIGHT_GRAY, sub_size=22
)

# Objection 2
add_text_slide(
    '"I already use billing software."',
    font_size=38, color=LIGHT_GRAY, bold=False
)
add_text_slide(
    "Software doesn't create discipline.\nStructure does.",
    font_size=42, color=GOLD,
    subtitle="Software records transactions.\nIt doesn't enforce them.",
    sub_color=LIGHT_GRAY, sub_size=22
)

# Objection 3
add_text_slide(
    '"My store is small."',
    font_size=38, color=LIGHT_GRAY, bold=False
)
add_text_slide(
    "Smaller stores leak MORE\nproportionally.",
    font_size=42, color=RED_WARM,
    subtitle="10% mismatch on ₹30L inventory = ₹3L uncontrolled stock.\nThat's not small.",
    sub_color=LIGHT_GRAY, sub_size=22
)

# Objection 4
add_text_slide(
    '"This sounds expensive."',
    font_size=38, color=LIGHT_GRAY, bold=False
)
add_text_slide(
    "Compare my fee to your\nmonthly hidden loss.",
    font_size=42, color=GOLD,
    subtitle="If you're losing ₹1–3 Lakh per year silently…\n₹85,000 pays for itself in 44 days.",
    sub_color=LIGHT_GRAY, sub_size=22
)

add_text_slide(
    "See how solvable this is?",
    font_size=48, color=TEAL
)


# ── SECTION 8: WHAT MAKES THIS DIFFERENT (Slides 81–88) ────────

add_text_slide(
    "The 30-Day Stock Certainty System™\nis different.",
    font_size=36, color=GOLD,
    accent_line=True, accent_color=GOLD
)

add_split_slide(
    "Others",
    [
        "❌  Give you software and leave",
        "❌  Sell features, not outcomes",
        "❌  No enforcement mechanism",
        "❌  No accountability",
        "❌  No measurable proof",
    ],
    title="OTHERS vs. US"
)

add_split_slide(
    "We",
    [
        "✅  Install structure + enforce discipline",
        "✅  Sell stock match, not software",
        "✅  30-day enforcement built-in",
        "✅  Weekly accountability reviews",
        "✅  Before-After validation report",
    ],
)

add_text_slide(
    "It's structured.\nIt's predictable.\nIt's measurable.\nIt's enforceable.",
    font_size=36, color=WHITE
)

add_text_slide(
    "We implement WITH you.\nAnd I guarantee results.",
    font_size=42, color=TEAL
)


# ── SECTION 9: GUARANTEE (Slides 89–92) ────────────────────────

add_text_slide(
    "The Stock Certainty\nGuarantee™",
    font_size=48, color=GOLD,
    accent_line=True, accent_color=GOLD
)

add_text_slide(
    "If mismatch doesn't reduce\nin 30 days…",
    font_size=42, color=WHITE
)

add_text_slide(
    "We work FREE\nuntil it does.",
    font_size=56, color=GREEN_OK,
    subtitle="Zero risk. I don't win unless you win.",
    sub_color=LIGHT_GRAY, sub_size=26
)


# ── SECTION 10: THE OFFER (Slides 93–102) ──────────────────────

add_text_slide(
    "Here's what you get:",
    font_size=44, color=GOLD,
    accent_line=True, accent_color=GOLD
)

add_split_slide(
    "The Complete\nStock Certainty\nSystem™",
    [
        "✅  Control Gap Audit™",
        "✅  Inventory Foundation Reset",
        "✅  Supplier Entry Lock™",
        "✅  Sales Deduction Lock™",
        "✅  Single System Enforcement",
        "✅  30-Day Compliance Monitoring",
        "✅  Before-After Validation Report",
    ],
    title="YOUR 30-DAY INSTALLATION"
)

add_split_slide(
    "BONUSES\nIncluded Free",
    [
        "🎁  Hidden Leakage Exposure Report™",
        "🎁  Staff Control Rulebook™",
        "🎁  10-Minute Daily Control Routine™",
        "🎁  90-Day Stability Audit™",
        "🎁  Expansion Readiness Scorecard™",
    ],
    title="+ HIGH-VALUE BONUSES",
    left_color=GOLD
)

add_number_slide(
    "₹85,000", "One-Time Investment",
    "Setup + 30-Day Installation + All Bonuses Included.\nNo hidden charges. No monthly lock-in for installation.",
    num_color=GOLD
)

add_text_slide(
    "Only 3 stores per month.",
    font_size=44, color=RED_WARM,
    subtitle="This is hands-on installation.\nNot bulk SaaS onboarding.\nSlots fill quickly.",
    sub_color=LIGHT_GRAY, sub_size=24,
    accent_line=True, accent_color=RED_WARM
)

add_text_slide(
    "Next cycle starts\non the 1st.",
    font_size=48, color=GOLD,
    subtitle="Miss this slot → wait 30 more days.\n30 more days of silent leakage.",
    sub_color=LIGHT_GRAY, sub_size=24
)


# ── SECTION 11: THE CTA (Slides 103–115) ───────────────────────

add_text_slide(
    "Want to see if this\nfits YOUR store?",
    font_size=44, color=WHITE,
    accent_line=True, accent_color=TEAL
)

add_text_slide(
    "Book a Free\nStock Leakage Audit.",
    font_size=48, color=TEAL,
    subtitle="30-minute detailed diagnosis of your store.\nWe identify your exact leakage zones.\nYou get clarity — instantly.",
    sub_color=LIGHT_GRAY, sub_size=24
)

add_text_slide(
    "No obligation.\nNo pressure.\nJust clarity for your store.",
    font_size=36, color=WHITE,
    subtitle="But only serious operators qualify.\nWe don't work with careless stores.",
    sub_color=LIGHT_GRAY, sub_size=22
)

add_text_slide(
    "Limited audit slots every week.",
    font_size=40, color=GOLD,
    subtitle="Once filled, the calendar closes.\nFirst serious store gets priority.",
    sub_color=LIGHT_GRAY, sub_size=24
)


# ── SECTION 12: FORK IN THE ROAD (Slides 116–130) ──────────────

add_text_slide(
    "You have two options.",
    font_size=48, color=WHITE
)

# Option A (bad)
add_text_slide(
    "Option A:",
    font_size=36, color=RED_WARM,
    subtitle="Ignore this.\nContinue monthly mismatch.\nContinue guessing profits.\nContinue the stress.",
    sub_color=LIGHT_GRAY, sub_size=28
)

# Option B (good)
add_text_slide(
    "Option B:",
    font_size=36, color=GREEN_OK,
    subtitle="Install a proven system.\nTake back control.\nProtect your margins.\nScale confidently.",
    sub_color=WHITE, sub_size=28
)

add_text_slide(
    "Makes complete\nbusiness sense.",
    font_size=48, color=GOLD
)

add_text_slide(
    "Your store deserves structure.\nYour hard work deserves protection.",
    font_size=36, color=WHITE
)

# Future vision
add_text_slide(
    "Imagine 6 months from now…",
    font_size=44, color=TEAL
)

add_text_slide(
    "Clean reports.\nConfident audits.\nGrowth plans.",
    font_size=44, color=WHITE,
    subtitle="But only if you act now.\nDoing nothing guarantees the same result.",
    sub_color=LIGHT_GRAY, sub_size=24
)


# ── SECTION 13: FINAL PUSH (Slides 131–138) ────────────────────

add_text_slide(
    "You're still watching.",
    font_size=48, color=WHITE,
    subtitle="That means you want control.\nSo act like it.",
    sub_color=GOLD, sub_size=30
)

add_text_slide(
    "📞 Book Your Free\nStock Leakage Audit",
    font_size=44, color=TEAL,
    subtitle="Takes less than 60 seconds.\nSlots fill quickly.\nNo commitment required.",
    sub_color=LIGHT_GRAY, sub_size=24,
    accent_line=True, accent_color=TEAL
)

add_text_slide(
    "Just clarity for your store.",
    font_size=44, color=WHITE,
    subtitle="We'll show you exactly where money is leaking\n— and how to stop it.",
    sub_color=GOLD, sub_size=26
)

add_text_slide(
    "Click the link below.\nBook your session.\nSecure your spot.",
    font_size=40, color=GOLD,
    accent_line=True, accent_color=GOLD
)

# Final slides
add_text_slide(
    "The 30-Day Stock Certainty System™\nworks.",
    font_size=40, color=WHITE,
    subtitle="We fix mismatch permanently.",
    sub_color=TEAL, sub_size=30
)

add_text_slide(
    "I'll see you inside.",
    font_size=52, color=GOLD,
    subtitle="— Aakash Savant\nRetail Control Architect™",
    sub_color=LIGHT_GRAY, sub_size=24
)


# ── Save ───────────────────────────────────────────────────────
out = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "VSL_Client_Demo_Presentation.pptx"
)
prs.save(out)
print(f"✅ Created: {out}")
print(f"Total slides: {len(prs.slides)}")
