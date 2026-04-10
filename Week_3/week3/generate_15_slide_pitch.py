"""
Generate a beautiful 15-slide Retail Control Architect pitch deck
for store owners - in simple, layman language (no jargons).
Designed for a 15-minute face-to-face presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette (Premium Dark Theme) ───
BG_DARK       = RGBColor(0x0F, 0x0F, 0x1A)   # Deep navy black
BG_CARD       = RGBColor(0x1A, 0x1A, 0x2E)   # Card background
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)   # Vibrant orange
ACCENT_GOLD   = RGBColor(0xFF, 0xD7, 0x00)   # Gold
ACCENT_GREEN  = RGBColor(0x00, 0xE6, 0x76)   # Success green
ACCENT_RED    = RGBColor(0xFF, 0x4D, 0x4D)   # Alert red
ACCENT_BLUE   = RGBColor(0x4D, 0xA8, 0xFF)   # Cool blue
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xB0, 0xB0, 0xC0)
MEDIUM_GRAY   = RGBColor(0x80, 0x80, 0x99)
SUBTLE_PURPLE = RGBColor(0x6C, 0x5C, 0xE7)   # Accent purple

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_solid_bg(slide, color):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_rect(slide, left, top, width, height, color1, color2):
    """Add a gradient rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2
    return shape


def add_accent_bar(slide, left, top, width, height, color):
    """Add a colored accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, color, opacity=1.0):
    """Add a decorative circle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=28,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.3):
    """Add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=24,
                       default_color=LIGHT_GRAY, font_name="Calibri", alignment=PP_ALIGN.LEFT):
    """
    Add multi-line text. Each line is a dict with keys:
    text, size (optional), color (optional), bold (optional), spacing_after (optional)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_info.get("text", "")
        p.font.size = Pt(line_info.get("size", default_size))
        p.font.color.rgb = line_info.get("color", default_color)
        p.font.bold = line_info.get("bold", False)
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(line_info.get("spacing_after", 8))
        p.space_before = Pt(line_info.get("spacing_before", 0))
    return txBox


def add_rounded_card(slide, left, top, width, height, fill_color=BG_CARD):
    """Add a card-like rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
    shape.line.width = Pt(1)
    return shape


def add_number_badge(slide, left, top, number, color=ACCENT_ORANGE):
    """Add a step number in a circle."""
    size = Inches(0.7)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return circle


def add_icon_shape(slide, left, top, icon_text, size=Inches(0.9), bg_color=ACCENT_ORANGE):
    """Add an emoji/icon in a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER
    return shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title / Hook
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_solid_bg(slide, BG_DARK)

# Decorative elements
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)
add_circle(slide, Inches(10.5), Inches(0.5), Inches(2.5), RGBColor(0xFF, 0x6B, 0x35))

# Title
add_text_box(slide, Inches(1), Inches(1.5), Inches(9), Inches(1.5),
             "Stop Losing Money", font_size=56, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.8), Inches(9), Inches(1.5),
             "From Your Own Store.", font_size=56, color=ACCENT_ORANGE, bold=True)

# Subtitle
add_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(1),
             "A simple 5-step system to find and fix hidden losses in your retail store — in just 30 days.",
             font_size=24, color=LIGHT_GRAY)

# Bottom bar
add_accent_bar(slide, Inches(1), Inches(6.2), Inches(2), Inches(0.06), ACCENT_ORANGE)
add_text_box(slide, Inches(1), Inches(6.4), Inches(6), Inches(0.5),
             "By Aakash Savant  •  Retail Control Architect™",
             font_size=16, color=MEDIUM_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: The Problem — "Does This Happen in Your Store?"
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_RED)

add_text_box(slide, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
             "Does This Happen In Your Store?", font_size=40, color=ACCENT_RED, bold=True)

problems = [
    ("🔴", "Stock count never matches what the computer shows"),
    ("🔴", "End of month — you're shocked by the numbers"),
    ("🔴", "Staff says 'I don't know where it went'"),
    ("🔴", "Items go missing but nobody can explain why"),
    ("🔴", "You feel stressed but can't pinpoint the exact problem"),
]

for i, (icon, text) in enumerate(problems):
    y = Inches(1.8) + Inches(i * 1.0)
    add_rounded_card(slide, Inches(1), y, Inches(10.5), Inches(0.8), BG_CARD)
    add_text_box(slide, Inches(1.3), y + Inches(0.12), Inches(0.6), Inches(0.6),
                 icon, font_size=28)
    add_text_box(slide, Inches(2.0), y + Inches(0.15), Inches(9), Inches(0.6),
                 text, font_size=22, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: The Real Cost — Leaking Bucket
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
             "Your Store is Like a Leaking Bucket 🪣", font_size=40, color=ACCENT_ORANGE, bold=True)

add_text_box(slide, Inches(1), Inches(1.7), Inches(10), Inches(0.8),
             "You keep pouring money in… but it keeps draining out silently.",
             font_size=24, color=LIGHT_GRAY)

leaks = [
    ("💸", "Sales happening but not recorded properly", ACCENT_RED),
    ("📦", "Returns taken back but never updated in stock", ACCENT_ORANGE),
    ("🚚", "Items moved between shelves / branches without writing it down", ACCENT_GOLD),
    ("🤷", "Small gaps every day = BIG losses every month", ACCENT_RED),
]

for i, (icon, text, color) in enumerate(leaks):
    y = Inches(2.8) + Inches(i * 1.05)
    add_accent_bar(slide, Inches(1), y, Inches(0.12), Inches(0.75), color)
    add_text_box(slide, Inches(1.5), y + Inches(0.1), Inches(0.6), Inches(0.6),
                 icon, font_size=28)
    add_text_box(slide, Inches(2.2), y + Inches(0.15), Inches(9), Inches(0.6),
                 text, font_size=22, color=WHITE)

# Bottom stat
add_rounded_card(slide, Inches(3.5), Inches(6.3), Inches(6), Inches(0.9), RGBColor(0x2A, 0x15, 0x15))
add_text_box(slide, Inches(3.7), Inches(6.4), Inches(5.5), Inches(0.7),
             "⚠️  Most stores lose ₹2-5 Lakhs every year without even knowing it!",
             font_size=18, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Real Story — ₹4 Lakh Gone  
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, SUBTLE_PURPLE)

add_text_box(slide, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
             "A Real Story From a Store Like Yours", font_size=40, color=WHITE, bold=True)

# Story cards
story_items = [
    ("Computer showed stock worth", "₹28 Lakhs", ACCENT_BLUE),
    ("Actual stock on shelves was only", "₹24 Lakhs", ACCENT_ORANGE),
    ("Money that silently disappeared", "₹4 Lakhs Gone!", ACCENT_RED),
]

for i, (label, value, color) in enumerate(story_items):
    y = Inches(1.9) + Inches(i * 1.5)
    add_rounded_card(slide, Inches(1.5), y, Inches(10), Inches(1.2), BG_CARD)
    add_text_box(slide, Inches(2), y + Inches(0.15), Inches(6), Inches(0.5),
                 label, font_size=20, color=LIGHT_GRAY)
    add_text_box(slide, Inches(2), y + Inches(0.6), Inches(6), Inches(0.5),
                 value, font_size=36, color=color, bold=True)

# Quote
add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.6),
             '"The owner thought this was normal. It\'s not."',
             font_size=20, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Root Cause — It's Not What You Think
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

add_text_box(slide, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
             "Why Does This Happen?", font_size=40, color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(1.6), Inches(10), Inches(0.8),
             "It's NOT what most store owners think…",
             font_size=24, color=LIGHT_GRAY)

# Wrong answers crossed out
wrong = [
    ("❌", "It's NOT a software problem", "Buying expensive software alone doesn't fix it"),
    ("❌", "It's NOT because of bad staff", "Blaming your team every day doesn't solve anything"),
]

right = [
    ("✅", "It's a SYSTEM problem", "There's no clear step-by-step process for daily stock control"),
]

for i, (icon, title, desc) in enumerate(wrong):
    y = Inches(2.6) + Inches(i * 1.3)
    add_rounded_card(slide, Inches(1), y, Inches(10.5), Inches(1.05), RGBColor(0x2A, 0x15, 0x15))
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(0.6), Inches(0.5),
                 icon, font_size=28)
    add_text_box(slide, Inches(2.0), y + Inches(0.08), Inches(8.5), Inches(0.5),
                 title, font_size=24, color=ACCENT_RED, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.55), Inches(8.5), Inches(0.5),
                 desc, font_size=18, color=MEDIUM_GRAY)

for i, (icon, title, desc) in enumerate(right):
    y = Inches(5.2) + Inches(i * 1.3)
    add_rounded_card(slide, Inches(1), y, Inches(10.5), Inches(1.05), RGBColor(0x0A, 0x2A, 0x15))
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(0.6), Inches(0.5),
                 icon, font_size=28)
    add_text_box(slide, Inches(2.0), y + Inches(0.08), Inches(8.5), Inches(0.5),
                 title, font_size=24, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.55), Inches(8.5), Inches(0.5),
                 desc, font_size=18, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Who Am I — Intro
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GOLD)

add_text_box(slide, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
             "Who Am I?", font_size=40, color=ACCENT_GOLD, bold=True)

# Name & role
add_text_box(slide, Inches(1), Inches(1.8), Inches(7), Inches(0.8),
             "Aakash Savant", font_size=44, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.7), Inches(7), Inches(0.5),
             "Founder — Retail Control Architect™", font_size=22, color=ACCENT_ORANGE)

# Credentials
creds = [
    "🎯  I help clothing & retail stores find and fix hidden stock losses",
    "🏪  Worked with multiple stores facing the same problem",
    "📊  Built a simple system that any store owner can follow",
    "💡  I focus on fixing the process, not just installing software",
]

for i, cred in enumerate(creds):
    y = Inches(3.6) + Inches(i * 0.85)
    add_text_box(slide, Inches(1.2), y, Inches(10), Inches(0.7),
                 cred, font_size=20, color=LIGHT_GRAY)

# Decorative circle
add_circle(slide, Inches(10.5), Inches(1.5), Inches(2), ACCENT_ORANGE)
add_text_box(slide, Inches(10.5), Inches(2.0), Inches(2), Inches(1),
             "AS", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Results / Proof
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GREEN)

add_text_box(slide, Inches(1), Inches(0.6), Inches(10), Inches(0.8),
             "Real Results From Real Stores", font_size=40, color=ACCENT_GREEN, bold=True)

results = [
    ("₹3.2 Lakhs", "Money recovered in just 45 days", "💰", ACCENT_GREEN),
    ("17% → 2%", "Stock mismatch reduced dramatically", "📉", ACCENT_BLUE),
    ("Zero Surprises", "No more month-end stock shocks", "✅", ACCENT_GOLD),
    ("2nd Branch", "Owner confidently opened another store", "🏪", ACCENT_ORANGE),
]

for i, (stat, desc, icon, color) in enumerate(results):
    x = Inches(0.8) + Inches(i * 3.1)
    add_rounded_card(slide, x, Inches(2.0), Inches(2.8), Inches(4.0), BG_CARD)
    
    # Icon circle
    icon_circle = add_circle(slide, x + Inches(0.85), Inches(2.4), Inches(1.0), color)
    add_text_box(slide, x + Inches(0.85), Inches(2.55), Inches(1.0), Inches(0.8),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.2), Inches(3.7), Inches(2.4), Inches(0.8),
                 stat, font_size=30, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(4.5), Inches(2.4), Inches(1.0),
                 desc, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: The Solution Overview
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
             "Introducing: Retail Control Architect™", font_size=40, color=ACCENT_ORANGE, bold=True)

add_text_box(slide, Inches(1), Inches(1.5), Inches(10), Inches(1.0),
             "A Simple 5-Step System to Stop Stock Losses",
             font_size=28, color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(2.3), Inches(10), Inches(0.8),
             "No complicated technology. No expensive software. Just a proven process that works.",
             font_size=20, color=LIGHT_GRAY)

# 5 steps preview
steps_preview = [
    ("1", "Find the Leaks", "🔍"),
    ("2", "Organize Your Stock", "📦"),
    ("3", "Hold Staff Accountable", "👥"),
    ("4", "Track Weekly", "📊"),
    ("5", "Lock the Process", "🔒"),
]

for i, (num, title, icon) in enumerate(steps_preview):
    x = Inches(0.5) + Inches(i * 2.5)
    add_rounded_card(slide, x, Inches(3.6), Inches(2.2), Inches(3.2), BG_CARD)
    
    # Step number
    add_number_badge(slide, x + Inches(0.75), Inches(3.9), num)
    
    # Icon
    add_text_box(slide, x + Inches(0.2), Inches(4.8), Inches(1.8), Inches(0.8),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    
    # Title
    add_text_box(slide, x + Inches(0.15), Inches(5.6), Inches(1.9), Inches(1.0),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Step 1 — Find the Leaks
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)

add_number_badge(slide, Inches(1), Inches(0.5), "1")
add_text_box(slide, Inches(2), Inches(0.5), Inches(8), Inches(0.8),
             "Find Where the Money is Leaking", font_size=36, color=WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(1.5), Inches(10), Inches(0.8),
             "We check your entire store — like a health checkup for your business.",
             font_size=22, color=LIGHT_GRAY)

items = [
    ("🔍", "Walk through your store and check real stock vs. computer stock"),
    ("📋", "Make a list of every place where items could go missing"),
    ("💡", "Find the TOP 3 reasons your store is losing money"),
    ("📊", "Give you a clear picture of how much you're actually losing"),
]

for i, (icon, text) in enumerate(items):
    y = Inches(2.6) + Inches(i * 1.1)
    add_rounded_card(slide, Inches(1), y, Inches(10.5), Inches(0.85), BG_CARD)
    add_text_box(slide, Inches(1.3), y + Inches(0.15), Inches(0.6), Inches(0.5),
                 icon, font_size=24)
    add_text_box(slide, Inches(2.0), y + Inches(0.18), Inches(9), Inches(0.5),
                 text, font_size=20, color=WHITE)

add_rounded_card(slide, Inches(3), Inches(6.4), Inches(7), Inches(0.7), RGBColor(0x0A, 0x2A, 0x15))
add_text_box(slide, Inches(3.2), Inches(6.48), Inches(6.5), Inches(0.5),
             "💡 Think of it as a doctor's diagnosis — before giving any medicine.",
             font_size=16, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Steps 2-3
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_BLUE)

# Step 2
add_number_badge(slide, Inches(0.8), Inches(0.5), "2", ACCENT_BLUE)
add_text_box(slide, Inches(1.7), Inches(0.5), Inches(5), Inches(0.7),
             "Organize Every Item Properly", font_size=30, color=WHITE, bold=True)

step2_items = [
    "Give every item a unique tag — like a name tag for each product",
    "Know exactly how many pieces of each item you have",
    "No more confusion about 'where did that item go?'",
]

for i, text in enumerate(step2_items):
    y = Inches(1.5) + Inches(i * 0.7)
    add_text_box(slide, Inches(1.2), y, Inches(10), Inches(0.6),
                 f"  •  {text}", font_size=18, color=LIGHT_GRAY)

# Divider
add_accent_bar(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.03), MEDIUM_GRAY)

# Step 3
add_number_badge(slide, Inches(0.8), Inches(4.1), "3", ACCENT_ORANGE)
add_text_box(slide, Inches(1.7), Inches(4.1), Inches(5), Inches(0.7),
             "Make Your Team Responsible", font_size=30, color=WHITE, bold=True)

step3_items = [
    "Each staff member knows their responsibilities clearly",
    "Simple daily checklist they must complete — takes only 5 minutes",
    "If something goes missing, you know exactly who to ask",
    "No more blame game or finger-pointing",
]

for i, text in enumerate(step3_items):
    y = Inches(5.0) + Inches(i * 0.6)
    add_text_box(slide, Inches(1.2), y, Inches(10), Inches(0.6),
                 f"  •  {text}", font_size=18, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: Steps 4-5
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GREEN)

# Step 4
add_number_badge(slide, Inches(0.8), Inches(0.5), "4", ACCENT_GREEN)
add_text_box(slide, Inches(1.7), Inches(0.5), Inches(5), Inches(0.7),
             "Check Progress Every Week", font_size=30, color=WHITE, bold=True)

step4_items = [
    "A simple one-page report shows how your store is doing",
    "See if stock is matching, if losses are reducing",
    "Catch problems EARLY — before they become big losses",
    "Takes just 10 minutes to review every week",
]

for i, text in enumerate(step4_items):
    y = Inches(1.5) + Inches(i * 0.7)
    add_text_box(slide, Inches(1.2), y, Inches(10), Inches(0.6),
                 f"  •  {text}", font_size=18, color=LIGHT_GRAY)

# Divider
add_accent_bar(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.03), MEDIUM_GRAY)

# Step 5
add_number_badge(slide, Inches(0.8), Inches(4.1), "5", ACCENT_GOLD)
add_text_box(slide, Inches(1.7), Inches(4.1), Inches(5), Inches(0.7),
             "Make It Permanent", font_size=30, color=WHITE, bold=True)

step5_items = [
    "Set clear rules for your store — everyone follows the same process",
    "Stock movement is always recorded — nothing moves without a record",
    "Your store runs smoothly even when you're not there",
    "System works on its own — no need to watch over every little thing",
]

for i, text in enumerate(step5_items):
    y = Inches(5.0) + Inches(i * 0.6)
    add_text_box(slide, Inches(1.2), y, Inches(10), Inches(0.6),
                 f"  •  {text}", font_size=18, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: Common Doubts / Objection Handling
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, SUBTLE_PURPLE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
             "You Might Be Thinking… 🤔", font_size=40, color=WHITE, bold=True)

doubts = [
    ("'My staff won't follow any system'",
     "→ That's exactly why our system has built-in checks. They can't skip steps."),
    ("'I already use billing software'",
     "→ Software alone doesn't create discipline. You need a PROCESS."),
    ("'My store is small, this won't apply'",
     "→ Smaller stores actually lose more percentage-wise. This works for any size."),
    ("'This sounds expensive'",
     "→ Compare our fee to the lakhs you lose silently every year."),
]

for i, (doubt, answer) in enumerate(doubts):
    y = Inches(1.6) + Inches(i * 1.4)
    add_rounded_card(slide, Inches(1), y, Inches(10.5), Inches(1.2), BG_CARD)
    add_text_box(slide, Inches(1.4), y + Inches(0.1), Inches(9.5), Inches(0.5),
                 doubt, font_size=20, color=ACCENT_RED, bold=True)
    add_text_box(slide, Inches(1.4), y + Inches(0.6), Inches(9.5), Inches(0.5),
                 answer, font_size=18, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13: Imagine Your Store 6 Months From Now
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GREEN)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
             "Imagine Your Store 6 Months From Now… ✨", font_size=40, color=ACCENT_GREEN, bold=True)

visions = [
    ("😌", "No more stress at month-end — numbers match every time", ACCENT_GREEN),
    ("📊", "You TRUST your reports — you know exactly where you stand", ACCENT_BLUE),
    ("💰", "You keep the money that was silently leaking out", ACCENT_GOLD),
    ("🏪", "Your store runs smoothly even on your off days", ACCENT_ORANGE),
    ("🚀", "You can confidently plan to grow — maybe open another branch", SUBTLE_PURPLE),
]

for i, (icon, text, color) in enumerate(visions):
    y = Inches(1.8) + Inches(i * 1.0)
    add_accent_bar(slide, Inches(1), y, Inches(0.1), Inches(0.75), color)
    add_text_box(slide, Inches(1.5), y + Inches(0.1), Inches(0.6), Inches(0.6),
                 icon, font_size=28)
    add_text_box(slide, Inches(2.2), y + Inches(0.15), Inches(9.5), Inches(0.6),
                 text, font_size=22, color=WHITE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14: Our Promise / Guarantee
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GOLD)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8),
             "Our Promise To You 🤝", font_size=40, color=ACCENT_GOLD, bold=True)

# Guarantee card
add_rounded_card(slide, Inches(2), Inches(1.8), Inches(9), Inches(2.0), RGBColor(0x15, 0x2A, 0x0A))
add_text_box(slide, Inches(2.5), Inches(2.0), Inches(8), Inches(0.6),
             "🛡️  ZERO-RISK GUARANTEE", font_size=28, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.5), Inches(2.7), Inches(8), Inches(1.0),
             "If your stock mismatch doesn't reduce in 30 days —\nwe'll keep working for FREE until it does. No questions asked.",
             font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

# What you get
whats = [
    ("🔍", "Free 30-minute Store Check-up — we find your biggest leak"),
    ("📋", "Clear Action Plan — exact steps to fix your store"),
    ("👨‍💼", "Personal Guidance — I work WITH you, not just give advice"),
    ("📞", "Ongoing Support — call me anytime during the program"),
]

for i, (icon, text) in enumerate(whats):
    y = Inches(4.3) + Inches(i * 0.75)
    add_text_box(slide, Inches(2.3), y, Inches(0.5), Inches(0.6), icon, font_size=22)
    add_text_box(slide, Inches(3.0), y + Inches(0.05), Inches(8), Inches(0.6),
                 text, font_size=19, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15: Call to Action
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, BG_DARK)
add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "You Have Two Choices Today…", font_size=40, color=WHITE, bold=True)

# Option A
add_rounded_card(slide, Inches(1), Inches(1.8), Inches(5.2), Inches(3.5), RGBColor(0x2A, 0x15, 0x15))
add_text_box(slide, Inches(1.3), Inches(2.0), Inches(4.5), Inches(0.6),
             "❌  Option A: Do Nothing", font_size=24, color=ACCENT_RED, bold=True)
nothing_items = [
    "Keep losing money silently every month",
    "Keep getting surprised at month-end",
    "Keep guessing your actual profits",
    "Keep stressing about stock problems",
]
for i, item in enumerate(nothing_items):
    y = Inches(2.8) + Inches(i * 0.6)
    add_text_box(slide, Inches(1.6), y, Inches(4.2), Inches(0.5),
                 f"•  {item}", font_size=16, color=MEDIUM_GRAY)

# Option B
add_rounded_card(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(3.5), RGBColor(0x0A, 0x2A, 0x15))
add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5), Inches(0.6),
             "✅  Option B: Take Control", font_size=24, color=ACCENT_GREEN, bold=True)
action_items = [
    "Stop all hidden stock losses in 30 days",
    "Know your real numbers with confidence",
    "Build a store that runs smoothly",
    "Grow your business fearlessly",
]
for i, item in enumerate(action_items):
    y = Inches(2.8) + Inches(i * 0.6)
    add_text_box(slide, Inches(7.4), y, Inches(4.8), Inches(0.5),
                 f"•  {item}", font_size=16, color=LIGHT_GRAY)

# CTA Button area
cta_shape = add_rounded_card(slide, Inches(3.5), Inches(5.8), Inches(6), Inches(1.2), ACCENT_ORANGE)
add_text_box(slide, Inches(3.5), Inches(5.9), Inches(6), Inches(0.6),
             "📞 Book Your FREE Store Check-up Today", font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.5), Inches(6.5), Inches(6), Inches(0.5),
             "Only a few slots available each week • No obligation",
             font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)


# ─── Save ───
output_path = "/Users/aakash/Desktop/Week_3/Retail_Control_15_Slides_Pitch.pptx"
prs.save(output_path)
print(f"✅ Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
