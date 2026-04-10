#!/usr/bin/env python3
"""
Generate a High-Impact, Steve Jobs Style Simple PPTX
For Aakash Savant — Retail Control Architect
Focus: Minimalist, Layman terms, Big Visual impact.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# --- Constants ---
# Colors
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(240, 189, 78)
GRAY = RGBColor(160, 160, 160)
BLUE = RGBColor(0, 122, 255) # Apple Blue
RED = RGBColor(255, 90, 90)

def add_slide_minimal(prs, major_text, minor_text=None, font_size=60, major_color=WHITE, bg_color=BLACK):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Major Text
    width = prs.slide_width
    height = prs.slide_height
    
    # Center Major Text
    txBox = slide.shapes.add_textbox(0, height/3, width, height/3)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = major_text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = major_color
    run.font.name = 'Helvetica' # Safe clean font

    # Minor Text (Subtitle)
    if minor_text:
        txBox_sub = slide.shapes.add_textbox(0, height/2 + Inches(0.5), width, height/4)
        tf_sub = txBox_sub.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.alignment = PP_ALIGN.CENTER
        run_sub = p_sub.add_run()
        run_sub.text = minor_text
        run_sub.font.size = Pt(28)
        run_sub.font.color.rgb = GRAY
        run_sub.font.name = 'Helvetica'

    return slide

def run():
    prs = Presentation()
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_slide_minimal(prs, "The Retail Control Authority", "The End of Chaos.", font_size=64, major_color=GOLD)

    # Slide 2: Problem - The Hole
    add_slide_minimal(prs, "A Hole in your Floor.", "Profit is leaking every single day.", font_size=72)

    # Slide 3: The Tax
    add_slide_minimal(prs, "₹2,739 / Day", "Your 'Chaos Tax' for unmanaged stock.", font_size=96, major_color=RED)

    # Slide 4: The Reality
    add_slide_minimal(prs, "Billing Software ≠ Control.", "You are still doing double the work.", font_size=54)

    # Slide 5: The Breakthrough
    add_slide_minimal(prs, "Introducing An Installation.", "Not a software. A physical structure.", font_size=60, major_color=BLUE)

    # Slide 6: The 10-Minute Pulse
    add_slide_minimal(prs, "Total Control in 10 Minutes.", "Check. Verify. Close. Go home.", font_size=60)

    # Slide 7: The Routine
    add_slide_minimal(prs, "Inward. Outward. Count.", "The three rules that save your store.", font_size=60, major_color=GOLD)

    # Slide 8: Result 1
    add_slide_minimal(prs, "System = Shop", "100% Accuracy. 100% of the time.", font_size=64)

    # Slide 9: Result 2
    add_slide_minimal(prs, "Ready for Store #2", "You finally have a foundation that scales.", font_size=64, major_color=BLUE)

    # Slide 10: The Guarantee
    add_slide_minimal(prs, "Stock Certainty Guarantee", "If it doesn't match, we work for free.", font_size=54, major_color=GOLD)

    # Slide 11: Final Call
    add_slide_minimal(prs, "Only 3 Slots for March.", "Let's stop the bleeding today.", font_size=60, major_color=WHITE)

    # Save
    out_file = "/Users/aakash/Desktop/Week_3/Aakash_Simpler_Control_Presentation.pptx"
    prs.save(out_file)
    print(f"PPTX created successfully at: {out_file}")

if __name__ == "__main__":
    run()
